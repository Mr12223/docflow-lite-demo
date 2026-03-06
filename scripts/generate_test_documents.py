"""
生成 DocFlow 测试样本文档。

输出格式：
- pdf（文本版 + 扫描版）
- docx / doc（兼容样本，doc 为 docx 改扩展名）
- xlsx / xls（兼容样本，xls 为 xlsx 改扩展名）
- pptx
- txt / md / json / csv
- png / jpg / bmp / tiff / webp
"""

from __future__ import annotations

import csv
import json
import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "sample_data" / "test_documents"

TITLE_ZH = "DocFlow 综合测试样本"
TITLE_EN = "DocFlow End-to-End Sample Pack"

TEXT_LINES_ZH = [
    "这是用于毕业设计联调的测试样本文档。",
    "内容覆盖中文、English、数字 12345、日期 2026-03-06。",
    "目标：验证 PDF、Word、Excel、PPT、文本、JSON、CSV 与图片 OCR。",
    "特殊字符：@ # % & * （）【】《》——。",
]

TEXT_LINES_EN = [
    "This file is generated for parser integration testing.",
    "It contains headings, tables, numbers, and mixed-language content.",
    "Expected outcome: stable extraction without mojibake.",
]

TABLE_ROWS = [
    ["编号", "姓名", "专业", "成绩", "备注"],
    ["001", "张三", "软件工程", "92", "正常"],
    ["002", "李四", "数据科学", "88", "包含英文 Notes"],
    ["003", "王五", "人工智能", "95", "含符号 #AI"],
]

JSON_DATA = {
    "project": "DocFlow",
    "purpose": "格式识别与内容提取测试",
    "date": "2026-03-06",
    "tags": ["pdf", "docx", "xlsx", "pptx", "ocr", "utf-8"],
    "summary": {
        "zh": "用于验证多格式文档提取结果是否正常。",
        "en": "Used to verify multi-format document extraction.",
    },
    "records": [
        {"id": 1, "name": "张三", "score": 92},
        {"id": 2, "name": "李四", "score": 88},
    ],
}


def ensure_output_dir() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def pick_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc") if bold else Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ]
    for font_path in candidates:
        if font_path.exists():
            return ImageFont.truetype(str(font_path), size=size)
    return ImageFont.load_default()


def write_text_files() -> None:
    txt_path = OUTPUT_DIR / "sample.txt"
    md_path = OUTPUT_DIR / "sample.md"
    json_path = OUTPUT_DIR / "sample.json"
    csv_path = OUTPUT_DIR / "sample.csv"

    txt_content = "\n".join(
        [TITLE_ZH, "=" * 24, *TEXT_LINES_ZH, "", *TEXT_LINES_EN, "", "表格预览："]
        + [" | ".join(row) for row in TABLE_ROWS]
    )
    txt_path.write_text(txt_content, encoding="utf-8")

    md_lines = [
        f"# {TITLE_ZH}",
        "",
        f"## {TITLE_EN}",
        "",
        *[f"- {line}" for line in TEXT_LINES_ZH],
        "",
        "### Sample Table",
        "",
        "| 编号 | 姓名 | 专业 | 成绩 | 备注 |",
        "| --- | --- | --- | --- | --- |",
    ]
    md_lines.extend([f"| {' | '.join(row)} |" for row in TABLE_ROWS[1:]])
    md_lines.extend(["", "### JSON Fields", "", "- project", "- purpose", "- records"])
    md_path.write_text("\n".join(md_lines), encoding="utf-8")

    json_path.write_text(json.dumps(JSON_DATA, ensure_ascii=False, indent=2), encoding="utf-8")

    with csv_path.open("w", encoding="utf-8", newline="") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerows(TABLE_ROWS)


def write_docx_and_doc() -> None:
    docx_path = OUTPUT_DIR / "sample.docx"
    doc_path = OUTPUT_DIR / "sample.doc"

    doc = Document()
    doc.add_heading(TITLE_ZH, level=0)
    doc.add_paragraph(TITLE_EN)

    for line in TEXT_LINES_ZH:
        doc.add_paragraph(line)

    doc.add_heading("English Section", level=1)
    for line in TEXT_LINES_EN:
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("测试表格", level=1)
    table = doc.add_table(rows=len(TABLE_ROWS), cols=len(TABLE_ROWS[0]))
    table.style = "Table Grid"
    for row_index, row in enumerate(TABLE_ROWS):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = value

    doc.add_heading("结论", level=1)
    doc.add_paragraph("如果解析结果正常，应能看到标题、段落、列表与表格。")
    doc.save(docx_path)

    shutil.copyfile(docx_path, doc_path)


def write_xlsx_and_xls() -> None:
    xlsx_path = OUTPUT_DIR / "sample.xlsx"
    xls_path = OUTPUT_DIR / "sample.xls"

    wb = Workbook()
    ws = wb.active
    ws.title = "概览"
    ws["A1"] = TITLE_ZH
    ws["A2"] = TITLE_EN
    ws["A4"] = "生成日期"
    ws["B4"] = "2026-03-06"
    ws["A5"] = "用途"
    ws["B5"] = "多格式提取测试"

    data_sheet = wb.create_sheet("成绩表")
    for row in TABLE_ROWS:
        data_sheet.append(row)

    notes_sheet = wb.create_sheet("说明")
    notes_sheet["A1"] = "说明"
    notes_sheet["A2"] = "该文件用于验证 Excel 解析、表格提取和数值统计。"
    notes_sheet["A3"] = "兼容样本 sample.xls 为同内容改扩展名版本。"

    wb.save(xlsx_path)
    shutil.copyfile(xlsx_path, xls_path)


def write_pptx() -> None:
    pptx_path = OUTPUT_DIR / "sample.pptx"
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = TITLE_ZH
    title_slide.placeholders[1].text = TITLE_EN

    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    bullet_slide.shapes.title.text = "测试内容"
    text_frame = bullet_slide.placeholders[1].text_frame
    text_frame.clear()
    for index, line in enumerate(TEXT_LINES_ZH + TEXT_LINES_EN[:2]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(20)

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "成绩表示例"
    rows = len(TABLE_ROWS)
    cols = len(TABLE_ROWS[0])
    table = table_slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.0), Inches(2.5)).table
    for row_index, row in enumerate(TABLE_ROWS):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = value

    prs.save(pptx_path)


def create_image(page_number: int) -> Image.Image:
    width, height = 1240, 1754
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)

    title_font = pick_font(48, bold=True)
    subtitle_font = pick_font(28, bold=True)
    body_font = pick_font(28)

    y = 90
    draw.text((80, y), f"{TITLE_ZH} - 第 {page_number} 页", fill="black", font=title_font)
    y += 90
    draw.text((80, y), TITLE_EN, fill="black", font=subtitle_font)
    y += 90

    lines = TEXT_LINES_ZH + TEXT_LINES_EN
    for line in lines:
        draw.text((80, y), line, fill="black", font=body_font)
        y += 58

    y += 20
    draw.text((80, y), "表格示例：", fill="black", font=subtitle_font)
    y += 70

    column_x = [80, 220, 430, 700, 860]
    row_height = 52
    for row_index, row in enumerate(TABLE_ROWS):
        top = y + row_index * row_height
        for col_index, value in enumerate(row):
            left = column_x[col_index]
            right = column_x[col_index + 1] if col_index + 1 < len(column_x) else 1140
            draw.rectangle([left, top, right, top + row_height], outline="black", width=2)
            draw.text((left + 10, top + 10), value, fill="black", font=body_font)

    footer_y = height - 120
    draw.line((80, footer_y, width - 80, footer_y), fill="gray", width=2)
    draw.text((80, footer_y + 25), f"Page {page_number} / OCR test image", fill="black", font=body_font)

    return image


def write_images_and_scan_pdf() -> None:
    image_page1 = create_image(1)
    image_page2 = create_image(2)

    image_page1.save(OUTPUT_DIR / "sample.png")
    image_page1.save(OUTPUT_DIR / "sample.jpg", quality=95)
    image_page1.save(OUTPUT_DIR / "sample.bmp")
    image_page1.save(OUTPUT_DIR / "sample.tiff")
    image_page1.save(OUTPUT_DIR / "sample.webp", quality=95)

    pdf_path = OUTPUT_DIR / "sample_scan.pdf"
    image_page1_rgb = image_page1.convert("RGB")
    image_page2_rgb = image_page2.convert("RGB")
    image_page1_rgb.save(pdf_path, save_all=True, append_images=[image_page2_rgb])


def pdf_escape(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def write_text_pdf() -> None:
    pdf_path = OUTPUT_DIR / "sample_text.pdf"
    lines = [
        "DocFlow ASCII PDF Sample",
        "This PDF is text-based and tests direct PDF text extraction.",
        "Date: 2026-03-06",
        "Fields: PDF / DOCX / XLSX / PPTX / TXT / CSV / JSON / OCR",
        "Status: Ready for parser smoke tests.",
    ]

    content_lines = ["BT", "/F1 14 Tf", "72 760 Td"]
    for index, line in enumerate(lines):
        if index > 0:
            content_lines.append("T*")
        content_lines.append(f"({pdf_escape(line)}) Tj")
    content_lines.append("ET")
    stream = "\n".join(content_lines).encode("ascii")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode("ascii") + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))

    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )

    pdf_path.write_bytes(pdf)


def write_manifest() -> None:
    manifest_path = OUTPUT_DIR / "README.md"
    manifest = [
        "# DocFlow 测试样本说明",
        "",
        "已生成以下文件：",
        "",
        "- `sample.txt`：纯文本样本",
        "- `sample.md`：Markdown 样本",
        "- `sample.json`：JSON 样本",
        "- `sample.csv`：CSV 表格样本",
        "- `sample.docx`：Word 样本",
        "- `sample.doc`：兼容样本（由 docx 改扩展名生成）",
        "- `sample.xlsx`：Excel 样本",
        "- `sample.xls`：兼容样本（由 xlsx 改扩展名生成）",
        "- `sample.pptx`：PPT 样本",
        "- `sample_text.pdf`：文本型 PDF 样本",
        "- `sample_scan.pdf`：扫描型 PDF 样本（适合测试 OCR 兜底）",
        "- `sample.png/.jpg/.bmp/.tiff/.webp`：图片 OCR 样本",
        "",
        "说明：",
        "",
        "- `sample.doc` 和 `sample.xls` 不是传统二进制老格式，而是兼容测试样本。",
        "- 当前项目对这两类文件采用“优先按新格式内容尝试打开”的策略。",
    ]
    manifest_path.write_text("\n".join(manifest), encoding="utf-8")


def main() -> None:
    ensure_output_dir()
    write_text_files()
    write_docx_and_doc()
    write_xlsx_and_xls()
    write_pptx()
    write_images_and_scan_pdf()
    write_text_pdf()
    write_manifest()
    print(f"测试样本已生成：{OUTPUT_DIR}")


if __name__ == "__main__":
    main()
