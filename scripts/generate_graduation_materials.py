# -*- coding: utf-8 -*-
from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt as PptPt


REPO_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = REPO_ROOT / "docs" / "graduation"
DOWNLOADS_DIR = Path(r"d:\Downloads")
TODAY = date.today()

PROJECT_META = {
    "topic": "DocFlow多格式文档处理工具设计与实现",
    "review_title": "多格式文档自动化处理与内容提取技术研究综述",
    "topic_source": "工程实践 / 自拟课题",
    "topic_type": "CX",
    "advisor": "待填写",
    "student_name": "待填写",
    "student_id": "待填写",
    "major": "软件工程",
    "class_name": "软件工程20-0X（待填写）",
    "college": "软件学院",
}

PROJECT_ANALYSIS = {
    "positioning": "DocFlow 是一个面向多格式文档内容提取的 Web 工具，目标是在统一界面内完成 PDF、Word、Excel、PPT、文本、JSON、CSV 与图片的解析、OCR 识别、结果导出与批量测试。",
    "architecture": [
        "前端采用单页 `HTML + CSS + JavaScript`，负责拖拽上传、任务轮询、结果展示、依赖检查和批量测试面板。",
        "后端采用 `Flask`，提供同步处理、异步处理、任务取消、依赖检测与安装、批量测试报告等接口。",
        "核心处理层位于 `docflow_core.py`，以解析器抽象统一封装 PDF、DOCX、XLSX、PPTX、TXT/MD/CSV/JSON 等格式。",
        "依赖与环境适配逻辑位于 `docflow_support.py`，用于生成依赖状态、错误分类、安装建议和降级路径。",
    ],
    "features": [
        "支持单文件处理、异步任务进度查询和取消。",
        "支持 OCR 兜底，当前代码兼容 RapidOCR、PaddleOCR、Tesseract、EasyOCR 多种引擎。",
        "支持 TXT、JSON、Markdown、CSV 四种输出格式，以及关键词、摘要、统计信息生成。",
        "支持批量测试、评测报告导出、依赖检查与一键安装建议，具备较强工程展示价值。",
        "支持 Docker 轻量部署，适合毕业设计演示和云端验收。",
    ],
    "test_summary": [
        "本地评测集 `reports/evaluation_20260306_172409/evaluation.json` 显示 6/6 规则校验通过，规则通过率 100%。",
        "批量测试报告 `reports/batch_test_20260306_203110/report.md` 涵盖 51 个样本，35/35 个带预期样本全部匹配。",
        "边界失败样本主要集中于空文件、损坏文件、伪造扩展名和不支持格式，说明项目已具备异常分类和容错设计。",
    ],
    "strengths": [
        "系统边界清晰，论文可以围绕“解析器调度 + OCR 兜底 + 工程部署 + 测试验证”展开。",
        "项目代码规模适中，既有算法实现，也有前后端联调、异常处理和工程部署内容。",
        "已有样本数据、批测脚本、报告目录和论文草稿，适合快速沉淀为开题材料。",
    ],
    "limitations": [
        "复杂扫描件与弱质量图像仍依赖 OCR 效果，准确率受字体、噪声和分辨率影响。",
        "旧版 `.doc/.xls` 仍需要 LibreOffice 等外部工具增强兼容性。",
        "当前是单机型演示系统，尚未引入数据库、鉴权、任务队列与持久化存储。",
    ],
}

OPENING_REPORT_SECTIONS = [
    (
        "一、调研资料的准备",
        "围绕“多格式文档自动化处理与内容提取”这一课题，前期已完成三类资料准备工作。第一，完成相关技术文献的查阅，重点关注文档版面分析、OCR 识别、文档图像理解、多模态文档模型、Python 文档处理库与轻量化 Web 部署等方向，形成了对国内外研究现状的初步认识。第二，完成项目代码与运行环境调研，已对当前仓库中的 `app.py`、`docflow_core.py`、`docflow_support.py`、`frontend/doc_tool.html`、`scripts/run_batch_tests.py` 等核心文件进行了分析，明确了系统采用 Flask 后端、前端单页页面、解析器抽象、OCR 兜底、批量测试与报告输出的总体实现路线。第三，完成样本与测试资料整理，项目内已包含常规样本、边界样本与历史测试报告，可为课题后续的功能验证、性能分析和论文实验章节提供基础数据。",
    ),
    (
        "二、设计的目的、要求、思路与预期成果",
        "本课题拟围绕当前项目进一步完成系统分析、方案固化、测试验证与论文写作，形成一套可演示、可复现、可答辩的毕业设计成果。设计目的在于解决 PDF、Word、Excel、PPT、文本、JSON、CSV 与图片等异构文档在统一场景下难以自动提取内容的问题，提高文档处理效率，降低人工整理成本。系统要求包括：一是实现多格式文件统一接入与自动识别；二是针对文本型与扫描型文档提供差异化处理路径；三是支持 OCR 兜底、关键词提取、摘要生成与多格式导出；四是具备批量测试、错误分类和依赖检查能力；五是界面友好、部署轻量、适合毕业设计现场展示。技术思路上，系统将采用“前端交互层—后端服务层—解析与分析层”的三层结构，在后端通过解析器工厂/调度机制统一封装 PDF、DOCX、XLSX、PPTX 与文本类文件的处理流程；对于扫描型 PDF 和图片，则采用 RapidOCR、Tesseract 等 OCR 引擎形成兜底链路；在结果层补充关键词、摘要、元数据与统计信息，并统一转换为 TXT、JSON、Markdown、CSV 等输出格式。预期成果包括：完成可运行的 DocFlow 系统，形成开题报告、文献综述、毕业论文、测试报告与答辩 PPT，并在论文中给出系统架构、关键模块、测试结果与改进方向。",
    ),
    (
        "三、任务完成的阶段内容及时间安排",
        "（1）2026年3月上旬—2026年3月中旬：完成开题准备、文献整理、需求分析与课题论证，输出开题报告、文献综述与开题答辩材料。\n（2）2026年3月中旬—2026年4月上旬：梳理现有代码结构，完成系统总体设计、模块划分、技术路线与数据流设计，补充系统架构图、处理流程图和用例分析。\n（3）2026年4月上旬—2026年4月下旬：围绕文档解析、OCR 兜底、输出格式化、前端交互和异常处理等模块进行功能完善与联调，补充样本测试与边界验证。\n（4）2026年4月下旬—2026年5月中旬：完成系统测试与结果统计，整理准确率、耗时、典型案例与问题分析，撰写论文初稿。\n（5）2026年5月中旬—2026年5月下旬：根据指导意见修改系统与论文内容，完善图表、参考文献和结论章节，形成论文定稿。\n（6）2026年5月下旬—2026年6月上旬：完成答辩 PPT、演示脚本与材料自检，准备毕业答辩。",
    ),
    (
        "四、完成设计（论文）所具备的条件因素",
        "本课题已具备较为完整的实施基础。首先，项目已有可运行的代码仓库、前后端页面、解析模块和批量测试脚本，降低了从零开发的风险。其次，开发环境与依赖较成熟，Python 生态内的 pdfplumber、PyMuPDF、python-docx、openpyxl、python-pptx、Flask、RapidOCR、Tesseract 等库能够满足主要功能需要。再次，项目内已有测试样本、评测报告和论文草稿，说明课题具有明确的工程对象与论文写作基础。最后，课题内容兼顾软件工程实现与一定的智能处理能力，既适合毕业设计要求，也具备进一步扩展为云端演示系统、任务队列系统和企业文档处理平台的空间。",
    ),
]

REVIEW_ABSTRACT = (
    "多格式文档处理是智能办公、档案数字化、知识管理与信息抽取中的基础环节。本文围绕 DocFlow 项目对应的“多格式文档自动化处理与内容提取”主题，从结构化文档解析、OCR 识别、版面理解、多模态文档建模以及轻量化工程部署等方面综述国内外研究进展。研究表明，pdfplumber、PyMuPDF、python-docx、openpyxl、python-pptx 等工具在文本型文档解析方面具有较高工程成熟度，而 Tesseract、RapidOCR 及 Donut、Nougat、LayoutLMv3、DocLLM 等方法则推动了扫描文档与复杂版面理解能力的持续提升。与此同时，通用大模型与端到端文档理解模型虽然增强了复杂页面分析和跨模态推理能力，但在算力成本、部署复杂度、响应时延和私有化场景适配方面仍存在较高门槛。结合本项目的实现路径，本文认为“规则解析 + OCR 兜底 + 统一调度 + 多格式导出 + 批量测试”是一条更适合本科毕业设计和中小型应用场景的技术路线。最后总结多格式文档处理的发展趋势，并给出对 DocFlow 项目设计与实现的启示。"
)

REVIEW_KEYWORDS = "文档处理；OCR；多模态文档理解；内容提取；DocFlow"

REVIEW_SECTIONS = [
    (
        "1 研究背景与问题定义",
        [
            "随着企事业单位业务数字化程度不断提升，合同、财务报表、技术文档、扫描档案、会议纪要和演示材料等信息大量以异构文档形式沉淀。如何从多格式文档中自动提取文本、表格、元数据与结构化信息，已经成为文档管理、知识工程和智能办公中的关键问题。BERT 等预训练语言模型推动了文本理解能力显著提升[1]，LayoutLM 系列将文本、版面和视觉特征联合建模，使文档理解从“纯文本处理”迈向“版面感知处理”[2-4]。这说明文档处理领域已从单一 OCR 或单一解析库阶段，发展到结构解析、视觉理解与多模态建模协同演进的新阶段。",
            "从工程视角看，多格式文档处理并不是简单的“读文件”。不同格式在内部组织方式、信息密度和错误模式上存在明显差异：PDF 既可能具有可提取文本层，也可能完全由扫描图像构成；DOCX、XLSX、PPTX 基于 OOXML 封装，结构稳定但样式与嵌入对象复杂；TXT、CSV、JSON 解析简单，但编码、空文件和非法内容问题频发；图片与扫描件则必须借助 OCR 才能获得文本信息。对于毕业设计项目而言，既要保证系统可运行、可展示和可部署，又要兼顾复杂样本的识别能力，因此技术路线必须在识别效果、工程复杂度与部署成本之间寻找平衡[13-17]。",
            "结合当前 DocFlow 仓库可以发现，项目已形成较清晰的问题定义：面向常见办公文档和图片输入，统一提供上传、解析、OCR 兜底、摘要与关键词生成、结果导出、批量测试和依赖检查能力。这种课题定位既契合软件工程专业对系统分析、模块设计、编码实现、测试验证和部署展示的要求，也为论文写作提供了可量化的实验对象。因而，文献综述不仅需要梳理相关算法和工具的发展脉络，更需要回答一个实际问题：在本科毕业设计的约束下，怎样构建一套轻量但完整、稳定且可扩展的文档处理系统。",
        ],
    ),
    (
        "2 多格式文档解析技术研究现状",
        [
            "在结构化或半结构化文档解析方面，业界通常优先采用文件格式对应的原生解析库。对于 PDF，pdfplumber 与 PyMuPDF 提供了对文本层、页面对象和部分布局元素的直接访问能力[8][12]；当 PDF 自身包含可检索文本时，这类方案具有速度快、字符保真度高、资源消耗低等优点。对于 DOCX、XLSX 和 PPTX，python-docx、openpyxl 与 python-pptx 已形成较成熟的 Python 工具链，能够较好地完成段落、表格、工作表、幻灯片文本和元数据读取[9-11]。这些方案在工程实践中具有非常高的可复用性，也是当前多数轻量型办公文档处理系统的首选基础能力。",
            "不过，基于规则和文件结构的解析也有天然边界。首先，PDF 的文本层质量并不总是可靠，常见问题包括字体映射异常、排版顺序错乱、多栏文本穿插和表格边界缺失。其次，Office 文档虽然结构清晰，但当文档中嵌入图片、复杂 SmartArt、批注或外部对象时，传统解析库往往只能覆盖主要文本内容。再次，旧版 `.doc/.xls` 文件与伪造扩展名文件容易触发兼容性或损坏问题，这要求系统在解析失败时具备清晰的错误分类和降级处理机制。文献与工程资料普遍表明，单一解析库并不能覆盖全部真实场景，因此“主解析 + 备用解析 + 错误兜底”的工程化设计更为稳妥[8-12]。",
            "近年来，文档处理研究进一步强调“版面理解”而非仅仅“字符提取”。LayoutLM、LayoutLMv2 与 LayoutLMv3 的连续发展说明，页面中的几何位置、视觉块关系和文本语义之间具有强耦合性[2-4]。国内相关研究也指出，复杂表单、票据和报告的关键信息提取，往往依赖版面分析、区域检测和结构建模的综合能力[18]。对于 DocFlow 这类以工程实现为主的系统而言，虽然不一定需要自研版面理解模型，但在系统设计中应为后续接入版面分析算法、表格结构恢复和字段抽取模块预留接口，这也是其论文创新点与后续扩展价值所在。",
        ],
    ),
    (
        "3 OCR 与多模态文档理解技术研究现状",
        [
            "OCR 是扫描文档处理的核心技术。传统 Tesseract 通过字符分割、语言模型与版面处理机制，在印刷体识别场景下长期保持较强实用价值，并因部署成熟、开源稳定而被广泛应用于工业系统[14]。随着深度学习的发展，CRAFFT 文本检测模型和 CRNN 文本识别模型显著提升了复杂场景文本检测与识别性能[20][21]。RapidOCR 等轻量型工具则在工程层面对检测、识别和推理框架进行整合，降低了 OCR 能力落地门槛[15]。从工程实践来看，轻量系统往往会保留 Tesseract 或 RapidOCR 作为默认 OCR 方案，在需要更高精度时再扩展更大模型或专用推理后端。",
            "与传统 OCR 不同，Donut 和 Nougat 等方法尝试直接完成“从页面图像到结构化文本”的端到端建模。Donut 提出 OCR-free 的文档理解思路，强调视觉编码与序列生成的统一[5]；Nougat 则聚焦学术文档场景，力图从复杂论文页面中直接恢复接近可编辑文本的结果[6]。这些方法体现出多模态生成模型在文档理解中的潜力，尤其适合具有复杂公式、图文混排和非规则布局的场景。然而，这类方法对训练数据、算力资源和模型部署环境要求更高，在中小型项目中直接落地仍存在成本与复杂度问题。",
            "进一步地，DocLLM 等研究将大语言模型与文档布局信息耦合，引入生成式、多模态、布局感知的统一建模框架[7]。这类工作提升了问答、信息抽取和跨区域推理能力，也代表了文档智能的发展方向。但就当前毕业设计课题而言，若过早引入大模型，会导致系统重心从“可运行的软件工程系统”转移到“数据与模型训练问题”，反而削弱项目的可控性和可答辩性。因此，在本科阶段更合理的路径是：以解析库和轻量 OCR 构成稳定主链路，将多模态模型作为对比研究或未来扩展方向写入论文，而不是全部功能都依赖大模型完成。",
            "综合现有研究可以发现，OCR 与文档理解正在从“字符识别”走向“结构理解”和“任务驱动理解”。这意味着未来系统评估指标也将不仅局限于字符准确率，而会进一步关注表格恢复质量、字段抽取正确率、摘要质量与问答表现。对于 DocFlow 项目来说，这一趋势启示我们：当前版本应优先把多格式解析、OCR 兜底、错误分类和统一输出做好；未来版本可在此基础上迭代版面分析、知识抽取、向量检索和文档问答功能，从而形成更完整的技术演进路线。",
        ],
    ),
    (
        "4 工程化实现方案对比与 DocFlow 的设计启示",
        [
            "从工程实现方式来看，当前文档处理系统大致可分为三类。第一类是云端文档 AI 平台，优点是识别能力强、服务完整，但往往依赖外部接口，部署成本和数据隐私约束较高。第二类是“重模型 + 服务化”方案，能够处理复杂任务，但需要 GPU、推理框架和更复杂的运维支持。第三类是“解析库 + OCR + Web 服务”的轻量工程方案，强调本地部署、离线运行、快速迭代和低门槛集成。对于本科毕业设计而言，第三类方案更适合作为主体实现，因为它更容易体现需求分析、模块划分、接口设计、异常处理、部署与测试等软件工程能力[13-17]。",
            "DocFlow 当前代码结构正体现了这种轻量工程化路线。前端通过单页页面完成文件拖拽、状态展示、批量测试、日志输出和结果下载；后端基于 Flask 提供同步处理与异步任务接口[13]；核心解析层按照文件类型组织解析器类，实现对 PDF、DOCX、XLSX、PPTX 和文本类文件的统一调度；支持模块则进一步封装依赖检测、错误信息归类和外部工具发现逻辑。与许多课程作业式系统相比，DocFlow 的优势在于它不只是“能跑通功能”，还已经形成了依赖面板、报告目录、批测脚本和 Docker 部署说明，这使其更接近一个具备演示价值的软件产品。",
            "从设计思想上看，DocFlow 的可借鉴之处主要体现在四个方面。其一，采用分层结构和解析器抽象，便于后续新增格式或替换 OCR 引擎。其二，区分文本型文档与扫描型文档的处理路径，避免把所有文件都交给 OCR，从而兼顾速度与准确率。其三，通过统一输出格式和统计字段，方便前端展示与结果下载，也有利于论文实验章节的数据整理。其四，引入批量测试和异常分类，使系统具备可验证性，而不仅是静态展示。上述设计与文献中强调的“可扩展、可维护、可评估”的工程目标是一致的[16][17][22]。",
            "当然，DocFlow 也存在进一步提升空间。首先，复杂表格结构恢复、图片内多区域识别、低质量扫描件处理仍有优化空间。其次，批量任务当前仍运行在 Web 进程内，面对大规模文件时可能需要队列、缓存和持久化机制支持。再次，结果层当前以提取、摘要和导出为主，尚未深入到实体识别、字段抽取、知识图谱构建或问答增强等更高层能力。对于毕业论文写作而言，这些不足并不是问题，反而可以作为系统局限性与未来工作章节的重要内容，使论文结构更加完整、论述更加真实。",
        ],
    ),
    (
        "5 总结",
        [
            "总体来看，多格式文档处理技术正在经历从格式解析、OCR 识别到多模态文档理解的演进。国外研究在版面感知预训练模型、端到端文档生成模型和大模型融合方面推进较快[2-7]，国内研究则在文档版面分析、智能办公与 OCR 工程落地方面持续积累[18][19]。但无论技术如何升级，实际系统建设仍需面对部署成本、响应时延、环境依赖、异常样本和可维护性等现实问题。对本科毕业设计而言，真正有价值的不是盲目追求最大模型，而是在有限条件下完成一套技术链条完整、工程逻辑清楚、测试结果可信的系统。",
            "因此，以 DocFlow 为对象开展毕业设计具有较强可行性和现实意义。该项目既能体现 Python 技术栈、文档处理工具链、OCR 兜底与前后端联调等工程实践能力，又能在论文中吸收文档理解、多模态建模与版面分析等研究成果作为理论支撑。基于现有项目基础，后续只需进一步梳理模块设计、补充实验数据、完善论文表达和答辩材料，就可以形成较完整的毕业设计成果。换言之，DocFlow 课题的核心价值不在于“模型多大”，而在于它将文献中的关键技术能力压缩为一个可部署、可测试、可展示的轻量系统，这正是本科软件工程毕业设计最需要突出的实践意义。",
        ],
    ),
]

REFERENCES = [
    "[1] Devlin J, Chang M W, Lee K, et al. BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding[C]//Proceedings of NAACL-HLT 2019. Minneapolis: ACL, 2019: 4171-4186.",
    "[2] Xu Y, Li M, Cui L, et al. LayoutLM: Pre-training of Text and Layout for Document Image Understanding[C]//Proceedings of KDD 2020. New York: ACM, 2020: 1192-1200.",
    "[3] Xu Y, Lv T, Cui L, et al. LayoutLMv2: Multi-modal Pre-training for Visually-Rich Document Understanding[C]//Proceedings of ACL-IJCNLP 2021. Online: ACL, 2021: 2579-2591.",
    "[4] Huang Y, Lv T, Cui L, et al. LayoutLMv3: Pre-training for Document AI with Unified Text and Image Masking[J/OL]. arXiv:2204.08387, 2022.",
    "[5] Kim G, Hong T, Yim M, et al. OCR-free Document Understanding Transformer[J/OL]. arXiv:2111.15664, 2022.",
    "[6] Blecher L, Cucurull G, Scialom T, et al. Nougat: Neural Optical Understanding for Academic Documents[J/OL]. arXiv:2308.13418, 2023.",
    "[7] Lyu C, Luo J, Huang T, et al. DocLLM: A Layout-Aware Generative Language Model for Multimodal Document Understanding[J/OL]. arXiv:2401.00908, 2024.",
    "[8] jsvine. pdfplumber Documentation[EB/OL]. [2026-03-10]. https://github.com/jsvine/pdfplumber.",
    "[9] python-openxml. python-docx Documentation[EB/OL]. [2026-03-10]. https://python-docx.readthedocs.io/.",
    "[10] openpyxl. openpyxl Documentation[EB/OL]. [2026-03-10]. https://openpyxl.readthedocs.io/.",
    "[11] scanny. python-pptx Documentation[EB/OL]. [2026-03-10]. https://python-pptx.readthedocs.io/.",
    "[12] Artifex. PyMuPDF Documentation[EB/OL]. [2026-03-10]. https://pymupdf.readthedocs.io/.",
    "[13] Pallets. Flask Documentation[EB/OL]. [2026-03-10]. https://flask.palletsprojects.com/.",
    "[14] Tesseract OCR. User Manual[EB/OL]. [2026-03-10]. https://tesseract-ocr.github.io/tessdoc/.",
    "[15] RapidAI. RapidOCR[EB/OL]. [2026-03-10]. https://github.com/RapidAI/RapidOCR.",
    "[16] McKinney W. Python for Data Analysis[M]. 3rd ed. Sebastopol: O'Reilly Media, 2022.",
    "[17] Géron A. Hands-On Machine Learning with Scikit-Learn, Keras, and TensorFlow[M]. 3rd ed. Sebastopol: O'Reilly Media, 2022.",
    "[18] 张子良, 陈文锋. 基于深度学习的文档版面分析方法研究[J]. 计算机工程与应用, 2022, 58(3): 112-120.",
    "[19] 王海峰, 吴华. 自然语言处理研究综述[J]. 中国计算机学会通讯, 2021, 17(10): 14-23.",
    "[20] Baek Y, Lee B, Han D, et al. Character Region Awareness for Text Detection[C]//Proceedings of CVPR 2019. Long Beach: IEEE, 2019: 9365-9374.",
    "[21] Shi B, Bai X, Yao C. An End-to-End Trainable Neural Network for Image-Based Sequence Recognition and Its Application to Scene Text Recognition[J]. IEEE Transactions on Pattern Analysis and Machine Intelligence, 2017, 39(11): 2298-2304.",
    "[22] 李航. 统计学习方法[M]. 2版. 北京: 清华大学出版社, 2019.",
]

PPT_SLIDES = [
    {
        "title": PROJECT_META["topic"],
        "body": [
            "开题答辩汇报",
            f"学生姓名：{PROJECT_META['student_name']}",
            f"专业班级：{PROJECT_META['class_name']}",
            f"指导教师：{PROJECT_META['advisor']}",
            f"日期：{TODAY.year}年{TODAY.month}月{TODAY.day}日",
        ],
        "type": "title",
    },
    {
        "title": "研究背景与意义",
        "body": [
            "办公与档案场景中存在大量 PDF、Word、Excel、PPT、图片等异构文档。",
            "人工提取内容效率低、重复劳动多，且容易出现漏项、错项与格式不统一问题。",
            "轻量化、可部署、可测试的文档处理平台具有明显的工程价值和展示价值。",
            "课题可同时体现软件工程能力与文档智能处理能力，适合作为毕业设计。 ",
        ],
    },
    {
        "title": "项目现状分析",
        "body": [
            "仓库已具备前后端分离原型、核心解析模块、依赖检查、批量测试和 Docker 部署说明。",
            "后端基于 Flask，支持同步处理、异步任务、进度轮询和任务取消。",
            "核心解析层支持 PDF、DOCX、XLSX、PPTX、TXT/MD/CSV/JSON 与图片 OCR。",
            "已有样本数据、评测报告和论文草稿，可直接转化为毕业设计材料。",
        ],
    },
    {
        "title": "系统总体架构",
        "body": [
            "前端交互层：文件上传、拖拽、日志、结果展示、批量测试、依赖面板。",
            "后端服务层：路由接口、任务调度、状态管理、报告服务、异常响应。",
            "解析与分析层：格式识别、解析器调度、OCR 兜底、关键词提取、摘要生成。",
            "输出与部署层：TXT/JSON/Markdown/CSV 导出、报告目录、Docker 演示部署。",
        ],
        "type": "architecture",
    },
    {
        "title": "研究内容与技术路线",
        "body": [
            "分析异构文档格式差异，构建统一处理入口与解析器抽象。",
            "针对文本型与扫描型文档设计“主解析 + OCR 兜底”双通路。",
            "完成关键词提取、摘要生成、多格式导出与错误分类。",
            "通过批量测试与边界样本验证系统功能、稳定性与可用性。",
        ],
    },
    {
        "title": "核心功能设计",
        "body": [
            "多格式文件上传与自动识别。",
            "PDF 文本提取与扫描 PDF OCR 兜底。",
            "Word / Excel / PPT / 文本类文件内容解析。",
            "图片 OCR、结果导出、依赖检测、一键批测与报告生成。",
        ],
    },
    {
        "title": "阶段成果与测试依据",
        "body": [
            "评测集规则校验 6/6 全部通过，规则通过率 100%。",
            "批量测试报告覆盖 51 个样本，35/35 个带预期样本全部匹配。",
            "系统已形成异常分类、错误提示和降级处理路径，适合论文实验章节整理。",
            "当前工作重点转向材料固化、论文撰写、答辩演示与局部优化。",
        ],
    },
    {
        "title": "拟解决的重点与难点",
        "body": [
            "复杂 PDF、图片和低质量扫描件的识别准确性仍需优化。",
            "旧版 `.doc/.xls` 兼容路径需要依赖外部工具增强。",
            "批量任务仍运行在 Web 进程内，后续可引入队列与持久化能力。",
            "论文中需平衡算法深度、工程完整性与可答辩性。",
        ],
    },
    {
        "title": "进度安排与预期成果",
        "body": [
            "3 月：完成开题、文综、总体设计与关键图表。",
            "4 月：完善功能实现、联调测试、收集实验数据。",
            "5 月：完成论文初稿、修改完善、形成定稿。",
            "6 月前：完成答辩 PPT、系统演示脚本和材料自检。",
            "预期成果：系统成品、论文、测试报告、开题材料与答辩 PPT。",
        ],
    },
    {
        "title": "总结",
        "body": [
            "DocFlow 课题具备真实项目基础、明确技术边界和较强工程可展示性。",
            "采用轻量解析库与 OCR 结合的路线，更适合本科毕业设计落地。",
            "当前已完成材料初稿生成，后续只需补充个人信息并按导师意见微调。",
            "谢谢老师审阅。",
        ],
    },
]


def find_template(patterns: list[str]) -> Path:
    for pattern in patterns:
        matches = sorted(DOWNLOADS_DIR.glob(pattern))
        if matches:
            return matches[0]
    raise FileNotFoundError(f"未找到模板文件：{patterns}")


def ensure_output_dir() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def get_rfonts(run) -> OxmlElement:
    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.append(r_fonts)
    return r_fonts


def style_run(run, size_pt: float, *, bold: bool = False, east_asia: str = "宋体", western: str = "Times New Roman") -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = western
    get_rfonts(run).set(qn("w:eastAsia"), east_asia)


def reset_paragraph(paragraph, text: str = ""):
    p = paragraph._element
    for child in list(p):
        if child.tag != qn("w:pPr"):
            p.remove(child)
    if text:
        run = paragraph.add_run(text)
        style_run(run, 12)
    return paragraph


def set_paragraph_text(paragraph, text: str, *, size_pt: float = 12, bold: bool = False, align: int | None = None, first_line_indent: float | None = None) -> None:
    reset_paragraph(paragraph)
    run = paragraph.add_run(text)
    style_run(run, size_pt, bold=bold)
    fmt = paragraph.paragraph_format
    fmt.line_spacing = 1.5
    fmt.space_before = Pt(0)
    fmt.space_after = Pt(0)
    if first_line_indent is not None:
        fmt.first_line_indent = Cm(first_line_indent)
    if align is not None:
        paragraph.alignment = align


def add_formatted_paragraph(document: Document, text: str, *, size_pt: float = 12, bold: bool = False, center: bool = False, first_line_indent: float | None = None) -> None:
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if center else WD_ALIGN_PARAGRAPH.JUSTIFY
    set_paragraph_text(paragraph, text, size_pt=size_pt, bold=bold, align=paragraph.alignment, first_line_indent=first_line_indent)


def remove_paragraph(paragraph) -> None:
    element = paragraph._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def clear_cell(cell) -> None:
    tc = cell._tc
    for child in list(tc):
        if child.tag != qn("w:tcPr"):
            tc.remove(child)
    cell.add_paragraph()


def add_cell_paragraph(cell, text: str, *, size_pt: float = 12, bold: bool = False, first_line_indent: float | None = None) -> None:
    paragraph = cell.paragraphs[-1]
    if paragraph.text:
        paragraph = cell.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    set_paragraph_text(paragraph, text, size_pt=size_pt, bold=bold, align=WD_ALIGN_PARAGRAPH.JUSTIFY, first_line_indent=first_line_indent)


def generate_analysis_markdown() -> Path:
    content = [
        "# DocFlow 项目分析",
        "",
        "## 1. 项目定位",
        PROJECT_ANALYSIS["positioning"],
        "",
        "## 2. 架构概览",
        *[f"- {item}" for item in PROJECT_ANALYSIS["architecture"]],
        "",
        "## 3. 主要功能",
        *[f"- {item}" for item in PROJECT_ANALYSIS["features"]],
        "",
        "## 4. 已有测试依据",
        *[f"- {item}" for item in PROJECT_ANALYSIS["test_summary"]],
        "",
        "## 5. 适合作为毕设的原因",
        *[f"- {item}" for item in PROJECT_ANALYSIS["strengths"]],
        "",
        "## 6. 当前局限与后续扩展",
        *[f"- {item}" for item in PROJECT_ANALYSIS["limitations"]],
        "",
        "## 7. 建议课题名称",
        f"- `{PROJECT_META['topic']}`",
        "- 如果学校更偏学术表述，也可使用“基于 Python 的多格式文档自动化处理与内容提取系统设计与实现”。",
        "",
        "## 8. 本次生成的材料",
        "- 开题报告表（基于用户提供模板生成）",
        "- 文献综述（基于用户提供模板生成）",
        "- 开题答辩 PPT",
        "",
        "## 9. 待补充信息",
        "- 学生姓名、学号、班级、指导教师姓名与职称未在仓库中明确给出，本次统一保留为“待填写”。",
        "- 若学校要求填写课题来源的具体类型，可在“工程实践 / 自拟课题”的基础上按导师意见修改。",
    ]
    output_path = OUTPUT_DIR / "项目分析.md"
    output_path.write_text("\n".join(content), encoding="utf-8")
    return output_path


def generate_opening_report_docx(template_path: Path) -> Path:
    doc = Document(str(template_path))
    table = doc.tables[0]
    table.cell(0, 1).text = PROJECT_META["topic"]
    table.cell(1, 1).text = PROJECT_META["topic_source"]
    table.cell(1, 3).text = PROJECT_META["topic_type"]
    table.cell(1, 5).text = PROJECT_META["advisor"]
    table.cell(2, 1).text = PROJECT_META["student_name"]
    table.cell(2, 3).text = PROJECT_META["student_id"]
    table.cell(2, 5).text = PROJECT_META["major"]

    for row_idx, col_idx in [(0, 1), (1, 1), (1, 3), (1, 5), (2, 1), (2, 3), (2, 5)]:
        cell = table.cell(row_idx, col_idx)
        for paragraph in cell.paragraphs:
            set_paragraph_text(paragraph, cell.text, size_pt=12, align=WD_ALIGN_PARAGRAPH.CENTER)

    content_cell = table.cell(3, 0)
    clear_cell(content_cell)
    add_cell_paragraph(content_cell, "开题报告内容：", size_pt=12, bold=True)
    for heading, paragraph_text in OPENING_REPORT_SECTIONS:
        add_cell_paragraph(content_cell, heading, size_pt=12, bold=True)
        for line in paragraph_text.split("\n"):
            add_cell_paragraph(content_cell, line, size_pt=12, first_line_indent=0.74)
    add_cell_paragraph(content_cell, f"指导教师签名：{PROJECT_META['advisor']}                          日期：{TODAY.year}年    月    日", size_pt=12)
    add_cell_paragraph(content_cell, "（本次生成稿保留个人信息待填写，可根据导师意见继续精简或扩展。）", size_pt=11)

    output_path = OUTPUT_DIR / "DocFlow_开题报告表_已生成.docx"
    doc.save(output_path)
    return output_path


def generate_literature_review_docx(template_path: Path) -> Path:
    doc = Document(str(template_path))

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text.startswith("题    目"):
            set_paragraph_text(paragraph, f"题    目               {PROJECT_META['review_title']}", size_pt=16, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("学生姓名"):
            set_paragraph_text(paragraph, f"学生姓名               {PROJECT_META['student_name']}", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("专业班级"):
            set_paragraph_text(paragraph, f"专业班级               {PROJECT_META['class_name']}", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("学    号"):
            set_paragraph_text(paragraph, f"学    号               {PROJECT_META['student_id']}", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("学    院"):
            set_paragraph_text(paragraph, f"学    院               {PROJECT_META['college']}", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("指导教师（职称）"):
            set_paragraph_text(paragraph, f"指导教师（职称）      {PROJECT_META['advisor']}（待填写）", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text == "***（***）":
            set_paragraph_text(paragraph, "", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)
        elif text.startswith("完成时间"):
            set_paragraph_text(paragraph, f"完成时间               {TODAY.year}年{TODAY.month}月{TODAY.day}日", size_pt=12, align=WD_ALIGN_PARAGRAPH.LEFT)

    body_start = 26
    for paragraph in list(doc.paragraphs[body_start:]):
        remove_paragraph(paragraph)

    add_formatted_paragraph(doc, PROJECT_META["review_title"], size_pt=16, bold=True, center=True)
    add_formatted_paragraph(doc, "")
    add_formatted_paragraph(doc, f"摘  要：{REVIEW_ABSTRACT}", size_pt=12, bold=False, center=False, first_line_indent=0.74)
    add_formatted_paragraph(doc, f"关键词：{REVIEW_KEYWORDS}", size_pt=12, bold=False, center=False)
    add_formatted_paragraph(doc, "")

    for title, paragraphs in REVIEW_SECTIONS:
        add_formatted_paragraph(doc, title, size_pt=16, bold=True, center=False)
        for paragraph in paragraphs:
            add_formatted_paragraph(doc, paragraph, size_pt=12, bold=False, center=False, first_line_indent=0.74)
        add_formatted_paragraph(doc, "")

    refs_heading = doc.add_paragraph()
    refs_heading.add_run().add_break(WD_BREAK.PAGE)
    add_formatted_paragraph(doc, "参考文献", size_pt=15, bold=True, center=True)
    add_formatted_paragraph(doc, "")
    for ref in REFERENCES:
        add_formatted_paragraph(doc, ref, size_pt=10.5, bold=False, center=False)

    output_path = OUTPUT_DIR / "DocFlow_文献综述_已生成.docx"
    doc.save(output_path)
    return output_path


def add_textbox(slide, left, top, width, height, text, *, font_size=24, bold=False, color=(31, 41, 55), align=PP_ALIGN.LEFT) -> None:
    text_box = slide.shapes.add_textbox(left, top, width, height)
    frame = text_box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Microsoft YaHei"


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(30, 64, 175)
    banner.line.color.rgb = RGBColor(30, 64, 175)
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(8.5), Inches(0.4), title, font_size=24, bold=True, color=(255, 255, 255))
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.6), Inches(5.8))
    frame = text_box.text_frame
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = PptPt(10)
        for run in paragraph.runs:
            run.font.size = PptPt(21)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = RGBColor(31, 41, 55)
    add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2.2), Inches(0.25), "DocFlow 开题答辩", font_size=10, color=(100, 116, 139), align=PP_ALIGN.RIGHT)


def add_title_slide(prs: Presentation, title: str, items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(239, 246, 255)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.8), Inches(11.8), Inches(1.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(30, 64, 175)
    banner.line.color.rgb = RGBColor(30, 64, 175)
    add_textbox(slide, Inches(1.1), Inches(1.15), Inches(11), Inches(0.8), title, font_size=28, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.2), Inches(2.6))
    frame = text_box.text_frame
    frame.word_wrap = True
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = PptPt(12)
        for run in paragraph.runs:
            run.font.size = PptPt(20 if idx == 0 else 18)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = RGBColor(30, 41, 59)
            if idx == 0:
                run.font.bold = True
    add_textbox(slide, Inches(3.9), Inches(6.6), Inches(5.3), Inches(0.3), "多格式文档自动化处理与内容提取", font_size=12, color=(71, 85, 105), align=PP_ALIGN.CENTER)


def add_architecture_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(30, 64, 175)
    banner.line.color.rgb = RGBColor(30, 64, 175)
    add_textbox(slide, Inches(0.5), Inches(0.1), Inches(8.5), Inches(0.4), title, font_size=24, bold=True, color=(255, 255, 255))

    cards = [
        ("前端层", "上传拖拽\n任务轮询\n结果展示\n批量测试", RGBColor(219, 234, 254)),
        ("服务层", "Flask 路由\n任务调度\n取消控制\n报告服务", RGBColor(224, 231, 255)),
        ("核心层", "解析器调度\nOCR 兜底\n关键词/摘要\n格式导出", RGBColor(220, 252, 231)),
        ("支撑层", "依赖检查\n错误分类\nDocker 部署\n样本与报告", RGBColor(254, 240, 138)),
    ]
    x_positions = [0.8, 3.55, 6.3, 9.05]
    for (name, body, color), x in zip(cards, x_positions):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.25), Inches(2.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(148, 163, 184)
        add_textbox(slide, Inches(x + 0.12), Inches(1.7), Inches(2.0), Inches(0.35), name, font_size=18, bold=True, color=(15, 23, 42), align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(2.15), Inches(1.95), Inches(1.45), body, font_size=15, color=(51, 65, 85), align=PP_ALIGN.CENTER)

    text_box = slide.shapes.add_textbox(Inches(0.85), Inches(4.4), Inches(11.1), Inches(2.1))
    frame = text_box.text_frame
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.space_after = PptPt(10)
        for run in paragraph.runs:
            run.font.size = PptPt(18)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = RGBColor(31, 41, 55)


def generate_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_spec in PPT_SLIDES:
        slide_type = slide_spec.get("type", "bullet")
        if slide_type == "title":
            add_title_slide(prs, slide_spec["title"], slide_spec["body"])
        elif slide_type == "architecture":
            add_architecture_slide(prs, slide_spec["title"], slide_spec["body"])
        else:
            add_bullet_slide(prs, slide_spec["title"], slide_spec["body"])

    output_path = OUTPUT_DIR / "DocFlow_开题答辩PPT_已生成.pptx"
    prs.save(output_path)
    return output_path


def generate_manifest(paths: list[Path]) -> Path:
    manifest = {
        "generated_at": TODAY.isoformat(),
        "project_topic": PROJECT_META["topic"],
        "files": [str(path.relative_to(REPO_ROOT)) for path in paths],
    }
    output_path = OUTPUT_DIR / "生成清单.json"
    output_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path


def main() -> None:
    ensure_output_dir()
    opening_template = find_template(["7-*.docx"])
    review_template = find_template(["8 *.docx"])

    outputs = [
        generate_analysis_markdown(),
        generate_opening_report_docx(opening_template),
        generate_literature_review_docx(review_template),
        generate_pptx(),
    ]
    outputs.append(generate_manifest(outputs))

    for item in outputs:
        print(item.relative_to(REPO_ROOT))


if __name__ == "__main__":
    main()
