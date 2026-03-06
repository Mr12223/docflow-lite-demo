"""
生成 DocFlow 的异常 / 边界测试样本。

输出目录：
- test_documents_edge_cases/

覆盖场景：
- 空文件 / 极短内容 / 超长内容
- 不同编码（UTF-8 BOM / GBK / UTF-16）
- 非法 JSON / 不规则 CSV
- 最小可用 Office 文件
- 伪造扩展名文件
- 空白 / 低对比度 / 超小字图片
- 损坏 PDF / 空白扫描 PDF
- 特殊文件名 / 不支持扩展名
"""

from __future__ import annotations

import csv
import json
import os
import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "sample_data" / "test_documents_edge_cases"


def ensure_output_dir() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def pick_font(size: int, bold: bool = False):
    candidates = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc") if bold else Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ]
    for font_path in candidates:
        if font_path.exists():
            return ImageFont.truetype(str(font_path), size=size)
    return ImageFont.load_default()


def write_text_edge_cases() -> None:
    (OUTPUT_DIR / "00_empty.txt").write_text("", encoding="utf-8")

    (OUTPUT_DIR / "01_utf8_bom.txt").write_text(
        "UTF-8 BOM 测试\n第一行：中文\nSecond line: English\n12345",
        encoding="utf-8-sig",
    )

    (OUTPUT_DIR / "02_gbk.txt").write_bytes(
        "GBK 编码测试：这是中文内容，用于验证自动编码检测。".encode("gbk")
    )

    (OUTPUT_DIR / "03_utf16.txt").write_text(
        "UTF-16 编码测试\n用于验证 TextParser 的编码识别能力。",
        encoding="utf-16",
    )

    long_line = "超长行测试|" + ("边界内容123ABC-" * 1500)
    (OUTPUT_DIR / "04_long_line.txt").write_text(long_line, encoding="utf-8")

    md_content = "\n".join(
        [
            "# Markdown 边界样本",
            "",
            "## 特殊字符",
            "",
            "- emoji: 😀 😄 🤖",
            "- symbols: !@#$%^&*()_+-=[]{};:'\",.<>/?\\|",
            "- mixed: 中文 English 12345",
            "",
            "```python",
            "def hello():",
            "    return 'markdown-edge-case'",
            "```",
        ]
    )
    (OUTPUT_DIR / "05_special_chars.md").write_text(md_content, encoding="utf-8")

    weird_name = OUTPUT_DIR / "06_带空格 和 中文 @2026!.txt"
    weird_name.write_text("特殊文件名测试。", encoding="utf-8")


def write_json_csv_edge_cases() -> None:
    valid_json = {
        "case": "nested",
        "items": [{"id": i, "value": f"项目-{i}"} for i in range(1, 6)],
        "meta": {"empty": "", "null": None, "bool": True},
    }
    (OUTPUT_DIR / "10_nested.json").write_text(
        json.dumps(valid_json, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    (OUTPUT_DIR / "11_invalid.json").write_text(
        '{\n  "name": "broken",\n  "value": 123,\n  "missing_end": true,\n',
        encoding="utf-8",
    )

    with (OUTPUT_DIR / "12_empty.csv").open("w", encoding="utf-8", newline="") as csvfile:
        csv.writer(csvfile).writerows([])

    with (OUTPUT_DIR / "13_irregular.csv").open("w", encoding="utf-8", newline="") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["编号", "姓名", "成绩"])
        writer.writerow(["001", "张三"])
        writer.writerow(["002", "李四", "88", "多出一列"])
        writer.writerow([])
        writer.writerow(["003", "王五", "95"])

    with (OUTPUT_DIR / "14_large_table.csv").open("w", encoding="utf-8", newline="") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["id", "name", "score", "note"])
        for i in range(1, 501):
            writer.writerow([i, f"学生{i}", 60 + i % 40, f"第{i}行数据"])


def write_word_edge_cases() -> None:
    empty_docx = OUTPUT_DIR / "20_empty.docx"
    doc = Document()
    doc.save(empty_docx)

    minimal_docx = OUTPUT_DIR / "21_minimal.docx"
    doc = Document()
    doc.add_heading("最小 Word 测试", level=1)
    doc.add_paragraph("只有一段正文。")
    doc.save(minimal_docx)

    compat_doc = OUTPUT_DIR / "22_compat.doc"
    shutil.copyfile(minimal_docx, compat_doc)

    fake_docx = OUTPUT_DIR / "23_fake.docx"
    fake_docx.write_text("这不是一个真正的 docx 文件。", encoding="utf-8")

    fake_doc = OUTPUT_DIR / "24_fake.doc"
    fake_doc.write_text("这不是一个真正的 doc 文件。", encoding="utf-8")


def write_excel_edge_cases() -> None:
    empty_xlsx = OUTPUT_DIR / "30_empty.xlsx"
    wb = Workbook()
    wb.save(empty_xlsx)

    minimal_xlsx = OUTPUT_DIR / "31_minimal.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "边界"
    ws["A1"] = "值"
    ws["A2"] = 0
    ws["A3"] = -1
    ws["A4"] = 3.14159
    ws["B2"] = "中文"
    wb.save(minimal_xlsx)

    compat_xls = OUTPUT_DIR / "32_compat.xls"
    shutil.copyfile(minimal_xlsx, compat_xls)

    fake_xlsx = OUTPUT_DIR / "33_fake.xlsx"
    fake_xlsx.write_text("not a real xlsx", encoding="utf-8")

    fake_xls = OUTPUT_DIR / "34_fake.xls"
    fake_xls.write_text("not a real xls", encoding="utf-8")


def write_ppt_edge_cases() -> None:
    empty_pptx = OUTPUT_DIR / "40_empty.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(empty_pptx)

    minimal_pptx = OUTPUT_DIR / "41_minimal.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "最小 PPT 测试"
    slide.placeholders[1].text = "只有一个标题页"
    prs.save(minimal_pptx)

    fake_pptx = OUTPUT_DIR / "42_fake.pptx"
    fake_pptx.write_text("not a real pptx", encoding="utf-8")


def create_text_image(
    path: Path,
    text_lines: list[str],
    *,
    size: tuple[int, int] = (1200, 900),
    background: str = "white",
    fill: str = "black",
    font_size: int = 36,
    spacing: int = 54,
    start_xy: tuple[int, int] = (80, 80),
    fmt: str | None = None,
) -> None:
    image = Image.new("RGB", size, background)
    draw = ImageDraw.Draw(image)
    font = pick_font(font_size)
    x, y = start_xy
    for line in text_lines:
        draw.text((x, y), line, fill=fill, font=font)
        y += spacing
    image.save(path, format=fmt)


def write_image_edge_cases() -> None:
    Image.new("RGB", (1200, 900), "white").save(OUTPUT_DIR / "50_blank.png")

    create_text_image(
        OUTPUT_DIR / "51_low_contrast.jpg",
        ["低对比度 OCR 测试", "灰底浅字，识别难度较高"],
        background="#d9d9d9",
        fill="#bfbfbf",
        fmt="JPEG",
    )

    create_text_image(
        OUTPUT_DIR / "52_tiny_text.png",
        ["超小字体测试", "tiny font", "12345"],
        font_size=12,
        spacing=20,
        start_xy=(20, 20),
        size=(320, 160),
    )

    create_text_image(
        OUTPUT_DIR / "53_dense_text.webp",
        [f"第{i:02d}行：用于测试密集文本 OCR 表现。ABCDE12345" for i in range(1, 21)],
        font_size=24,
        spacing=32,
        size=(1200, 900),
        fmt="WEBP",
    )

    create_text_image(
        OUTPUT_DIR / "54_tiff_scan.tiff",
        ["TIFF 扫描样本", "用于测试图片 OCR 与格式兼容性"],
        fmt="TIFF",
    )


def build_blank_scan_pdf() -> None:
    page1 = Image.new("RGB", (1240, 1754), "white")
    page2 = Image.new("RGB", (1240, 1754), "#f6f6f6")
    pdf_path = OUTPUT_DIR / "60_blank_scan.pdf"
    page1.save(pdf_path, save_all=True, append_images=[page2])


def build_corrupt_pdf() -> None:
    content = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >>\n"
    )
    (OUTPUT_DIR / "61_corrupt.pdf").write_bytes(content)


def build_text_pdf_edge() -> None:
    pdf_path = OUTPUT_DIR / "62_short_text.pdf"
    stream = b"BT\n/F1 12 Tf\n72 720 Td\n(Short PDF edge test.) Tj\nET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode("ascii") + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode("ascii"))
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    pdf_path.write_bytes(pdf)


def write_misc_cases() -> None:
    (OUTPUT_DIR / "70_unsupported.bin").write_bytes(os.urandom(256))
    (OUTPUT_DIR / "71_empty.json").write_text("", encoding="utf-8")


def write_readme() -> None:
    lines = [
        "# DocFlow 异常 / 边界测试样本",
        "",
        "## 文件与预期",
        "",
        "| 文件 | 说明 | 预期表现 |",
        "| --- | --- | --- |",
        "| `00_empty.txt` | 空文本文件 | 成功，提取 0 字符 |",
        "| `01_utf8_bom.txt` | UTF-8 BOM 编码 | 成功，中文正常 |",
        "| `02_gbk.txt` | GBK 编码文本 | 成功，自动识别编码 |",
        "| `03_utf16.txt` | UTF-16 文本 | 成功，自动识别编码 |",
        "| `04_long_line.txt` | 超长单行文本 | 成功，测试长文本稳定性 |",
        "| `05_special_chars.md` | 特殊字符 Markdown | 成功，保留内容 |",
        "| `06_带空格 和 中文 @2026!.txt` | 特殊文件名 | 成功，验证文件名兼容 |",
        "| `10_nested.json` | 嵌套 JSON | 成功，格式化输出 |",
        "| `11_invalid.json` | 非法 JSON | 失败，返回解析错误 |",
        "| `12_empty.csv` | 空 CSV | 成功，行列为 0 |",
        "| `13_irregular.csv` | 列数不一致 CSV | 成功，按原样读取 |",
        "| `14_large_table.csv` | 500 行 CSV | 成功，测试大表格 |",
        "| `20_empty.docx` | 空 Word 文档 | 成功或接近空内容 |",
        "| `21_minimal.docx` | 最小 Word 文档 | 成功 |",
        "| `22_compat.doc` | 兼容 .doc 样本 | 成功 |",
        "| `23_fake.docx` | 伪造 docx | 失败，返回错误 |",
        "| `24_fake.doc` | 伪造 doc | 失败，返回错误 |",
        "| `30_empty.xlsx` | 空 Excel | 成功 |",
        "| `31_minimal.xlsx` | 最小 Excel | 成功 |",
        "| `32_compat.xls` | 兼容 .xls 样本 | 成功 |",
        "| `33_fake.xlsx` | 伪造 xlsx | 失败，返回错误 |",
        "| `34_fake.xls` | 伪造 xls | 失败，返回错误 |",
        "| `40_empty.pptx` | 空白 PPT | 成功或近空内容 |",
        "| `41_minimal.pptx` | 最小 PPT | 成功 |",
        "| `42_fake.pptx` | 伪造 pptx | 失败，返回错误 |",
        "| `50_blank.png` | 空白图片 | OCR 结果为空或极少 |",
        "| `51_low_contrast.jpg` | 低对比度图片 | OCR 可能不稳定 |",
        "| `52_tiny_text.png` | 超小字体图片 | OCR 可能部分缺失 |",
        "| `53_dense_text.webp` | 密集文本图片 | OCR 可提取部分内容 |",
        "| `54_tiff_scan.tiff` | TIFF 样本 | OCR 成功 |",
        "| `60_blank_scan.pdf` | 空白扫描 PDF | 失败，提示未识别到内容 |",
        "| `61_corrupt.pdf` | 损坏 PDF | 失败，返回错误 |",
        "| `62_short_text.pdf` | 极短文本 PDF | 成功 |",
        "| `70_unsupported.bin` | 不支持扩展名 | 失败，提示格式不支持 |",
        "| `71_empty.json` | 空 JSON 文件 | 失败，返回解析错误 |",
        "",
        "## 说明",
        "",
        "- `22_compat.doc` 与 `32_compat.xls` 是兼容测试样本，不是真正旧版二进制格式。",
        "- 这些样本适合用于演示“系统在异常输入下的鲁棒性与错误提示”。",
    ]
    (OUTPUT_DIR / "README.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    ensure_output_dir()
    write_text_edge_cases()
    write_json_csv_edge_cases()
    write_word_edge_cases()
    write_excel_edge_cases()
    write_ppt_edge_cases()
    write_image_edge_cases()
    build_blank_scan_pdf()
    build_corrupt_pdf()
    build_text_pdf_edge()
    write_misc_cases()
    write_readme()
    print(f"异常 / 边界样本已生成：{OUTPUT_DIR}")


if __name__ == "__main__":
    main()
