"""
DocFlow - 多格式文档自动化处理与内容提取工具
毕业设计核心代码
"""

import os
import re
import json
import time
import logging
import sys
import importlib
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

from docflow_support import prepare_pytesseract

# 修复 Windows 编码问题
if sys.platform == 'win32':
    import io
    if hasattr(sys.stdout, 'buffer') and not isinstance(sys.stdout, io.TextIOWrapper):
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    if hasattr(sys.stderr, 'buffer') and not isinstance(sys.stderr, io.TextIOWrapper):
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger("DocFlow")


class DocFlowCancelledError(Exception):
    pass


def _ensure_not_cancelled(cancel_callback=None) -> None:
    if not callable(cancel_callback):
        return
    try:
        cancelled = bool(cancel_callback())
    except DocFlowCancelledError:
        raise
    except Exception:
        return
    if cancelled:
        raise DocFlowCancelledError("任务已取消")


def _get_base_site_packages() -> Optional[Path]:
    base_prefix = getattr(sys, 'base_prefix', sys.prefix)
    candidates = []

    if os.name == 'nt':
        candidates.append(Path(base_prefix) / 'Lib' / 'site-packages')

    candidates.extend([
        Path(base_prefix) / 'lib' / f'python{sys.version_info.major}.{sys.version_info.minor}' / 'site-packages',
        Path(base_prefix) / 'lib' / 'site-packages',
    ])

    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def import_with_base_fallback(module_name: str):
    try:
        return importlib.import_module(module_name)
    except ModuleNotFoundError as exc:
        candidate = _get_base_site_packages()
        if not candidate:
            raise exc

        candidate_str = str(candidate)
        if candidate_str not in sys.path:
            sys.path.append(candidate_str)

        return importlib.import_module(module_name)


# ═══════════════════════════════════════════════════
#  数据结构定义
# ═══════════════════════════════════════════════════

@dataclass
class ExtractionResult:
    """文档提取结果数据类"""
    file_path: str
    file_type: str
    text_content: str = ""
    tables: list = field(default_factory=list)       # List[List[List[str]]]
    images: list = field(default_factory=list)        # List[str] — base64 or path
    metadata: dict = field(default_factory=dict)
    statistics: dict = field(default_factory=dict)
    processing_time_ms: float = 0.0
    success: bool = True
    error_msg: str = ""

    def to_dict(self):
        return {
            "file_path": self.file_path,
            "file_type": self.file_type,
            "text_content": self.text_content,
            "tables": self.tables,
            "metadata": self.metadata,
            "statistics": self.statistics,
            "processing_time_ms": self.processing_time_ms,
            "success": self.success,
            "error_msg": self.error_msg,
        }


# ═══════════════════════════════════════════════════
#  解析器基类
# ═══════════════════════════════════════════════════

class BaseParser:
    """所有解析器的抽象基类"""
    supported_extensions: tuple = ()

    def can_parse(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower().lstrip('.')
        return ext in self.supported_extensions

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        raise NotImplementedError

    def _make_result(self, file_path: str) -> ExtractionResult:
        return ExtractionResult(
            file_path=file_path,
            file_type=Path(file_path).suffix.lower().lstrip('.')
        )


# ═══════════════════════════════════════════════════
#  PDF 解析器
# ═══════════════════════════════════════════════════

class PDFParser(BaseParser):
    """
    PDF 文档解析器
    优先逐页使用 pdfplumber / PyMuPDF 提取文本；
    对提取不到文本的页面，再使用高分辨率渲染 + OCR + 预处理增强。
    """
    supported_extensions = ('pdf',)
    OCR_RENDER_SCALES = (2.0, 1.6)
    OCR_BINARY_THRESHOLDS = (180,)
    MIN_TEXT_SIGNAL = 8

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        start = time.time()
        result = self._make_result(file_path)
        progress_callback = kwargs.get("progress_callback")
        cancel_callback = kwargs.get("cancel_callback")

        def report(progress_pct: float, stage: str, message: str = "", **extra) -> None:
            if callable(progress_callback):
                progress_callback(
                    progress_pct=max(0.0, min(float(progress_pct), 100.0)),
                    stage=stage,
                    message=message,
                    **extra,
                )

        try:
            _ensure_not_cancelled(cancel_callback)
            pdf_mode = self._resolve_pdf_mode(kwargs.get("pdf_mode"))
            report(4, "pdf_prepare", "正在读取 PDF 结构")
            text_content, tables, metadata = self._parse_mixed(
                file_path,
                pdf_mode=pdf_mode,
                progress_callback=report,
                cancel_callback=cancel_callback,
            )
            report(92, "pdf_finalize", "正在汇总页级结果")
            result.text_content = text_content
            result.tables = tables
            result.metadata = metadata
        except DocFlowCancelledError:
            raise
        except Exception as e:
            result.error_msg = str(e)
            result.success = False
            report(100, "error", str(e))

        result.processing_time_ms = (time.time() - start) * 1000
        result.statistics = self._compute_stats(result.text_content, result.tables)
        logger.info(f"PDF解析完成: {file_path} — {result.statistics.get('char_count',0)} 字符")
        return result

    def _resolve_pdf_mode(self, value: Optional[str]) -> str:
        return value if value in {"accurate", "balanced", "fast"} else "balanced"

    def _get_pdf_mode_config(self, pdf_mode: str) -> dict:
        configs = {
            "accurate": {
                "render_scales": (2.2, 1.8),
                "binary_thresholds": (170, 190),
                "min_text_signal": 6,
                "target_long_edge": 1500,
                "max_long_edge": 1800,
                "ocr_variants": ("rgb", "gray", "contrast", "binary"),
                "score_break_threshold": 120,
                "variant_break_threshold": 95,
                "prefer_fast_ocr": False,
                "label": "高精度",
            },
            "balanced": {
                "render_scales": (2.0, 1.6),
                "binary_thresholds": (180,),
                "min_text_signal": 8,
                "target_long_edge": 1200,
                "max_long_edge": 1600,
                "ocr_variants": ("rgb", "gray", "contrast", "binary"),
                "score_break_threshold": 80,
                "variant_break_threshold": 70,
                "prefer_fast_ocr": False,
                "label": "平衡",
            },
            "fast": {
                "render_scales": (1.4,),
                "binary_thresholds": (),
                "min_text_signal": 12,
                "target_long_edge": 1000,
                "max_long_edge": 1200,
                "ocr_variants": ("rgb", "gray"),
                "score_break_threshold": 60,
                "variant_break_threshold": 48,
                "prefer_fast_ocr": True,
                "label": "快速",
            },
        }
        return configs[pdf_mode]

    def _build_pdf_runtime_config(
        self,
        file_path: str,
        pdf_mode: str,
        page_count: int,
        has_pdfplumber: bool,
        has_pymupdf: bool,
    ) -> dict:
        config = dict(self._get_pdf_mode_config(pdf_mode))
        file_size_mb = 0.0
        try:
            file_size_mb = Path(file_path).stat().st_size / (1024 * 1024)
        except Exception:
            pass

        large_by_pages = page_count >= 40
        large_by_size = file_size_mb >= 18
        huge_pdf = page_count >= 120 or file_size_mb >= 60
        is_large_pdf = large_by_pages or large_by_size

        config.update(
            {
                "file_size_mb": round(file_size_mb, 2),
                "is_large_pdf": is_large_pdf,
                "is_huge_pdf": huge_pdf,
                "prefer_text_engine": "pdfplumber" if has_pdfplumber else ("pymupdf" if has_pymupdf else "ocr"),
                "extract_tables": has_pdfplumber,
                "table_page_cap": None,
                "runtime_label": config["label"],
            }
        )

        if has_pymupdf and (pdf_mode == "fast" or is_large_pdf):
            config["prefer_text_engine"] = "pymupdf"

        if pdf_mode == "fast" and has_pdfplumber:
            config["extract_tables"] = False
            config["table_page_cap"] = 0

        if is_large_pdf:
            config["prefer_fast_ocr"] = True
            if pdf_mode == "balanced":
                config["render_scales"] = (1.6,)
                config["binary_thresholds"] = ()
                config["ocr_variants"] = ("gray", "rgb")
                config["target_long_edge"] = min(config["target_long_edge"], 1100)
                config["max_long_edge"] = min(config["max_long_edge"], 1350)
                config["variant_break_threshold"] = min(config["variant_break_threshold"], 58)
                config["score_break_threshold"] = min(config["score_break_threshold"], 72)
                config["table_page_cap"] = 18 if has_pdfplumber else 0
                config["runtime_label"] = "大文件优化 / 平衡"
            elif pdf_mode == "fast":
                config["render_scales"] = (1.2,)
                config["ocr_variants"] = ("gray",)
                config["target_long_edge"] = min(config["target_long_edge"], 920)
                config["max_long_edge"] = min(config["max_long_edge"], 1100)
                config["variant_break_threshold"] = min(config["variant_break_threshold"], 42)
                config["score_break_threshold"] = min(config["score_break_threshold"], 55)
                config["table_page_cap"] = 0
                config["runtime_label"] = "大文件优化 / 快速"
            else:
                config["table_page_cap"] = 24 if has_pdfplumber else 0
                config["runtime_label"] = "大文件优化 / 高精度"

        if huge_pdf:
            config["prefer_text_engine"] = "pymupdf" if has_pymupdf else config["prefer_text_engine"]
            if pdf_mode != "accurate":
                config["extract_tables"] = False
                config["table_page_cap"] = 0
            config["runtime_label"] = f"{config['runtime_label']}（超大文档）"

        return config

    def _should_extract_tables(self, page_index: int, mode_config: dict) -> bool:
        if not mode_config.get("extract_tables", True):
            return False
        table_page_cap = mode_config.get("table_page_cap")
        if table_page_cap is None:
            return True
        return page_index < int(table_page_cap)

    def _parse_mixed(self, file_path: str, pdf_mode: str = "balanced", progress_callback=None, cancel_callback=None) -> tuple[str, list, dict]:
        mode_config = self._get_pdf_mode_config(pdf_mode)
        text_parts = []
        all_tables = []
        page_sources = []
        open_errors = {}
        pdfplumber_doc = None
        pymupdf_doc = None
        pdfium_doc = None
        pdf_info = {}
        page_count = 0
        last_ocr_error = None
        runtime_text_engine = ""
        ocr_disabled_reason = ""

        def emit(progress_pct: float, stage: str, message: str = "", **extra) -> None:
            if callable(progress_callback):
                progress_callback(
                    progress_pct=max(0.0, min(float(progress_pct), 100.0)),
                    stage=stage,
                    message=message,
                    **extra,
                )

        try:
            _ensure_not_cancelled(cancel_callback)
            pdfplumber_doc = self._open_pdfplumber(file_path)
            page_count = len(pdfplumber_doc.pages)
            pdf_info = pdfplumber_doc.metadata or {}
            primary_engine = "pdfplumber"
            try:
                pymupdf_doc = self._open_pymupdf(file_path)
                if not pdf_info:
                    pdf_info = pymupdf_doc.metadata or {}
            except Exception as fitz_exc:
                open_errors["pymupdf"] = str(fitz_exc)
        except Exception as exc:
            open_errors["pdfplumber"] = str(exc)
            logger.warning(f"pdfplumber 打开失败，尝试 PyMuPDF: {exc}")
            try:
                pymupdf_doc = self._open_pymupdf(file_path)
                page_count = pymupdf_doc.page_count
                pdf_info = pymupdf_doc.metadata or {}
                primary_engine = "pymupdf"
            except Exception as fitz_exc:
                open_errors["pymupdf"] = str(fitz_exc)
                logger.warning(f"PyMuPDF 打开失败，切换 OCR 全页兜底: {fitz_exc}")
                pdfium_doc = self._open_pdfium(file_path)
                page_count = len(pdfium_doc)
                primary_engine = "ocr"

        mode_config = self._build_pdf_runtime_config(
            file_path=file_path,
            pdf_mode=pdf_mode,
            page_count=page_count,
            has_pdfplumber=pdfplumber_doc is not None,
            has_pymupdf=pymupdf_doc is not None,
        )
        runtime_text_engine = mode_config.get("prefer_text_engine", primary_engine)
        if primary_engine != "ocr":
            primary_engine = runtime_text_engine
        emit(8, "pdf_open", f"已识别 {page_count} 页，准备逐页提取")
        if mode_config.get("is_large_pdf"):
            logger.info(
                "PDF 大文件优化已启用: %s ｜ pages=%s ｜ size=%sMB ｜ text=%s ｜ tables=%s",
                mode_config.get("runtime_label"),
                page_count,
                mode_config.get("file_size_mb", 0),
                primary_engine,
                "on" if mode_config.get("extract_tables") else "off",
            )

        try:
            for index in range(page_count):
                _ensure_not_cancelled(cancel_callback)
                page_text = ""
                page_tables = []
                page_source = primary_engine
                page_error = ""
                progress_start = 10 + (index / max(page_count, 1)) * 78
                progress_end = 10 + ((index + 1) / max(page_count, 1)) * 78

                def emit_page(ratio: float, stage: str, message: str = "", **extra) -> None:
                    emit(
                        progress_start + (progress_end - progress_start) * max(0.0, min(float(ratio), 1.0)),
                        stage,
                        message,
                        current_page=index + 1,
                        total_pages=page_count,
                        **extra,
                    )

                emit_page(0.03, "pdf_page_prepare", f"第 {index + 1}/{page_count} 页：读取文本层")

                if primary_engine == "pdfplumber" and pdfplumber_doc is not None:
                    try:
                        page = pdfplumber_doc.pages[index]
                        page_text = self._normalize_pdf_text(page.extract_text() or "")
                    except Exception as exc:
                        page_text = ""
                        page_error = f"pdfplumber: {exc}"
                elif primary_engine == "pymupdf" and pymupdf_doc is not None:
                    try:
                        page_text = self._normalize_pdf_text(
                            self._extract_pymupdf_text(pymupdf_doc, index)
                        )
                    except Exception as exc:
                        page_text = ""
                        page_error = f"pymupdf: {exc}"

                if not self._has_meaningful_text(page_text, min_signal=mode_config["min_text_signal"]):
                    if primary_engine == "pymupdf" and pdfplumber_doc is not None:
                        try:
                            alt_page = pdfplumber_doc.pages[index]
                            alt_text = self._normalize_pdf_text(alt_page.extract_text() or "")
                            if self._has_meaningful_text(alt_text, min_signal=mode_config["min_text_signal"]):
                                page_text = alt_text
                                page_source = "pdfplumber"
                                page_error = ""
                        except Exception:
                            pass
                    elif primary_engine == "pdfplumber" and pymupdf_doc is not None:
                        try:
                            alt_text = self._normalize_pdf_text(
                                self._extract_pymupdf_text(pymupdf_doc, index)
                            )
                            if self._has_meaningful_text(alt_text, min_signal=mode_config["min_text_signal"]):
                                page_text = alt_text
                                page_source = "pymupdf"
                                page_error = ""
                        except Exception:
                            pass

                if not self._has_meaningful_text(page_text, min_signal=mode_config["min_text_signal"]) and not ocr_disabled_reason:
                    try:
                        _ensure_not_cancelled(cancel_callback)
                        if pdfium_doc is None:
                            pdfium_doc = self._open_pdfium(file_path)
                        emit_page(0.42, "pdf_page_ocr", f"第 {index + 1}/{page_count} 页：文本较弱，启动 OCR")
                        ocr_text = self._ocr_pdf_page(
                            pdfium_doc,
                            index,
                            mode_config=mode_config,
                            progress_callback=emit_page,
                            page_no=index + 1,
                            total_pages=page_count,
                            cancel_callback=cancel_callback,
                        )
                        if self._has_meaningful_text(ocr_text, min_signal=mode_config["min_text_signal"]):
                            page_text = ocr_text
                            page_source = "ocr"
                    except Exception as exc:
                        last_ocr_error = exc
                        exc_text = str(exc)
                        if any(key in exc_text for key in ("无法初始化 OCR 引擎", "Tesseract 初始化失败", "未检测到 Tesseract 可执行程序", "未安装 pytesseract")):
                            ocr_disabled_reason = exc_text
                            logger.warning("OCR 引擎初始化失败，后续页面将跳过 OCR：%s", exc_text)
                        if not page_error:
                            page_error = f"ocr: {exc}"
                elif not self._has_meaningful_text(page_text, min_signal=mode_config["min_text_signal"]) and ocr_disabled_reason and not page_error:
                    page_error = f"ocr: {ocr_disabled_reason}"

                if pdfplumber_doc is not None and self._should_extract_tables(index, mode_config):
                    try:
                        table_page = pdfplumber_doc.pages[index]
                        page_tables = self._extract_pdfplumber_tables(table_page)
                        if page_tables:
                            all_tables.extend(page_tables)
                    except Exception as exc:
                        if not page_error:
                            page_error = f"table: {exc}"

                if page_text:
                    text_parts.append(f"[第 {index + 1} 页]\n{page_text}")

                page_sources.append(
                    {
                        "page": index + 1,
                        "source": page_source,
                        "chars": len(page_text),
                        "tables": len(page_tables),
                        "error": page_error,
                    }
                )
                page_info = "文本层" if page_source != "ocr" else "OCR"
                if not page_text:
                    page_info = "未提取到有效内容"
                emit_page(
                    1.0,
                    "pdf_page_done",
                    f"第 {index + 1}/{page_count} 页完成：{page_info}",
                    page_source=page_source,
                    chars=len(page_text),
                    tables=len(page_tables),
                    page_error=page_error,
                )
        finally:
            try:
                if pdfplumber_doc is not None:
                    pdfplumber_doc.close()
            except Exception:
                pass
            try:
                if pymupdf_doc is not None:
                    pymupdf_doc.close()
            except Exception:
                pass
            try:
                if pdfium_doc is not None:
                    pdfium_doc.close()
            except Exception:
                pass

        text_content = "\n\n".join(text_parts).strip()
        if not text_content:
            if last_ocr_error is not None:
                raise RuntimeError(str(last_ocr_error))

            if open_errors:
                failure_summary = "；".join(f"{name}: {msg}" for name, msg in open_errors.items())
                raise RuntimeError(f"PDF 解析失败：无法读取文本层，也未识别到可用 OCR 内容。{failure_summary}")

        metadata = {
            "pages": page_count,
            "info": pdf_info,
            "text_pages": sum(1 for item in page_sources if item["chars"] > 0 and item["source"] != "ocr"),
            "ocr_pages": sum(1 for item in page_sources if item["chars"] > 0 and item["source"] == "ocr"),
            "table_pages": sum(1 for item in page_sources if item["tables"] > 0),
            "page_sources": page_sources,
            "primary_engine": primary_engine,
            "ocr_engine": getattr(self, "_active_ocr_engine", "EasyOCR / Tesseract"),
            "pdf_mode": pdf_mode,
            "pdf_mode_label": mode_config["label"],
            "pdf_runtime_label": mode_config.get("runtime_label", mode_config["label"]),
            "file_size_mb": mode_config.get("file_size_mb", 0),
            "large_pdf_optimized": bool(mode_config.get("is_large_pdf")),
            "table_extraction_enabled": bool(mode_config.get("extract_tables")),
            "table_page_cap": mode_config.get("table_page_cap"),
            "table_extraction": (
                "已关闭（大文件优化）"
                if not mode_config.get("extract_tables")
                else (
                    f"前 {mode_config['table_page_cap']} 页尝试提取表格"
                    if mode_config.get("table_page_cap") not in (None, 0)
                    else "全页尝试提取表格"
                )
            ),
        }
        if open_errors:
            metadata["fallback_reason"] = open_errors
        if metadata["ocr_pages"] > 0 and metadata["text_pages"] > 0:
            metadata["解析方式"] = "逐页混合提取（文本层 + OCR）"
        elif metadata["ocr_pages"] > 0:
            metadata["解析方式"] = "高分辨率 OCR"
        else:
            metadata["解析方式"] = primary_engine
        return text_content, all_tables, metadata

    def _open_pdfplumber(self, file_path: str):
        pdfplumber = import_with_base_fallback('pdfplumber')
        return pdfplumber.open(file_path)

    def _open_pymupdf(self, file_path: str):
        fitz = import_with_base_fallback('fitz')
        return fitz.open(file_path)

    def _open_pdfium(self, file_path: str):
        pdfium = import_with_base_fallback('pypdfium2')
        return pdfium.PdfDocument(file_path)

    def _extract_pymupdf_text(self, doc, page_index: int) -> str:
        page = doc.load_page(page_index)
        return page.get_text("text") or ""

    def _extract_pdfplumber_tables(self, page) -> list:
        tables = []
        try:
            for table in page.extract_tables():
                if not table:
                    continue
                clean_table = [
                    ["" if cell is None else str(cell).strip() for cell in row]
                    for row in table
                    if row
                ]
                if clean_table:
                    tables.append(clean_table)
        except Exception as exc:
            logger.warning(f"pdfplumber 表格提取失败: {exc}")
        return tables

    def _normalize_pdf_text(self, text: str) -> str:
        text = (text or "").replace("\x00", " ")
        lines = [line.rstrip() for line in text.splitlines()]
        normalized = "\n".join(lines)
        normalized = re.sub(r"\n{3,}", "\n\n", normalized)
        return normalized.strip()

    def _has_meaningful_text(self, text: str, min_signal: Optional[int] = None) -> bool:
        cleaned = self._normalize_pdf_text(text)
        if not cleaned:
            return False
        signal = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", cleaned))
        threshold = self.MIN_TEXT_SIGNAL if min_signal is None else min_signal
        return signal >= threshold or len(cleaned) >= 20

    def _ocr_pdf_page(self, pdf, page_index: int, mode_config: Optional[dict] = None, progress_callback=None, page_no: Optional[int] = None, total_pages: Optional[int] = None, cancel_callback=None) -> str:
        import tempfile

        mode_config = mode_config or self._get_pdf_mode_config("balanced")
        _ensure_not_cancelled(cancel_callback)
        backend = self._get_ocr_backend(prefer_fast=bool(mode_config.get("prefer_fast_ocr")))
        page = pdf[page_index]
        best_text = ""
        best_score = float("-inf")
        last_error = None

        try:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                temp_image_path = tmp_file.name

            try:
                total_scales = max(len(mode_config["render_scales"]), 1)
                for scale_index, scale in enumerate(mode_config["render_scales"], start=1):
                    _ensure_not_cancelled(cancel_callback)
                    if callable(progress_callback):
                        progress_callback(
                            0.48 + ((scale_index - 1) / total_scales) * 0.22,
                            "pdf_ocr_render",
                            f"第 {page_no or page_index + 1}/{total_pages or '?'} 页：OCR 渲染 x{scale}",
                        )
                    image = page.render(scale=scale).to_pil().convert("RGB")
                    try:
                        page_text = self._ocr_pil_image(
                            backend,
                            image,
                            temp_image_path,
                            mode_config=mode_config,
                            progress_callback=(
                                (lambda ratio, stage, message:
                                    progress_callback(
                                        0.55 + (((scale_index - 1) + max(0.0, min(float(ratio), 1.0))) / total_scales) * 0.35,
                                        stage,
                                        message,
                                    )
                                ) if callable(progress_callback) else None
                            ),
                            page_no=page_no or page_index + 1,
                            total_pages=total_pages,
                            cancel_callback=cancel_callback,
                        )
                    except Exception as exc:
                        last_error = exc
                        logger.warning(f"PDF 第 {page_index + 1} 页 OCR 失败（scale={scale}）: {exc}")
                        continue
                    finally:
                        try:
                            image.close()
                        except Exception:
                            pass

                    score = self._score_ocr_text(page_text)
                    if page_text and score > best_score:
                        best_text = page_text
                        best_score = score
                    if best_score >= mode_config["score_break_threshold"]:
                        break
            finally:
                try:
                    os.remove(temp_image_path)
                except Exception:
                    pass
        finally:
            try:
                page.close()
            except Exception:
                pass

        if best_text:
            return best_text
        if last_error is not None:
            raise last_error
        return ""

    def _ocr_pil_image(self, backend: dict, image, temp_image_path: str, mode_config: Optional[dict] = None, progress_callback=None, page_no: Optional[int] = None, total_pages: Optional[int] = None, cancel_callback=None) -> str:
        mode_config = mode_config or self._get_pdf_mode_config("balanced")
        best_text = ""
        best_score = float("-inf")
        last_error = None

        candidates = list(self._iter_ocr_candidates(image, mode_config=mode_config))
        total_candidates = max(len(candidates), 1)
        for candidate_index, (variant_name, candidate) in enumerate(candidates, start=1):
            try:
                _ensure_not_cancelled(cancel_callback)
                if callable(progress_callback):
                    progress_callback(
                        (candidate_index - 1) / total_candidates,
                        "pdf_ocr_variant",
                        f"第 {page_no or '?'} / {total_pages or '?'} 页：尝试 {variant_name} 方案",
                    )
                if backend["kind"] == "easyocr":
                    candidate.save(temp_image_path, format='PNG')
                    text = "\n".join(backend["engine"].readtext(temp_image_path, detail=0)).strip()
                else:
                    text = backend["engine"].image_to_string(
                        candidate,
                        lang=backend.get("lang", "chi_sim+eng"),
                        config=backend.get("config", ""),
                    ).strip()
                score = self._score_ocr_text(text)
                if text and score > best_score:
                    best_text = text
                    best_score = score
                if best_score >= mode_config.get("variant_break_threshold", float("inf")):
                    break
            except Exception as exc:
                last_error = exc
                logger.warning(f"{backend['name']} 识别失败，预处理方案 {variant_name}: {exc}")

        if best_text:
            return self._normalize_pdf_text(best_text)
        if last_error is not None:
            raise last_error
        return ""

    def _iter_ocr_candidates(self, image, mode_config: Optional[dict] = None):
        from PIL import ImageEnhance, ImageOps

        mode_config = mode_config or self._get_pdf_mode_config("balanced")
        enabled = set(mode_config.get("ocr_variants") or ())
        base = image.convert("RGB")
        if "rgb" in enabled:
            yield "rgb", self._fit_image_for_ocr(base, mode_config=mode_config)

        gray = ImageOps.grayscale(base)
        gray = ImageOps.autocontrast(gray)
        if "gray" in enabled:
            yield "gray_autocontrast", self._fit_image_for_ocr(gray, mode_config=mode_config)

        contrast = ImageEnhance.Contrast(gray).enhance(1.8)
        sharpened = ImageEnhance.Sharpness(contrast).enhance(2.2)
        if "contrast" in enabled:
            yield "contrast_sharp", self._fit_image_for_ocr(sharpened, mode_config=mode_config)

        if "binary" in enabled:
            for threshold in mode_config.get("binary_thresholds") or ():
                binary = sharpened.point(lambda px, t=threshold: 255 if px > t else 0, mode="1").convert("L")
                yield f"binary_{threshold}", self._fit_image_for_ocr(binary, mode_config=mode_config)

    def _fit_image_for_ocr(self, image, mode_config: Optional[dict] = None):
        from PIL import Image

        mode_config = mode_config or self._get_pdf_mode_config("balanced")
        width, height = image.size
        long_edge = max(width, height)
        max_long_edge = int(mode_config.get("max_long_edge", 1600))
        target_long_edge = int(mode_config.get("target_long_edge", 1200))
        if long_edge > max_long_edge:
            scale = max_long_edge / long_edge
            return image.resize((max(1, int(width * scale)), max(1, int(height * scale))), Image.LANCZOS)
        if long_edge >= target_long_edge:
            return image

        scale = target_long_edge / max(long_edge, 1)
        return image.resize((max(1, int(width * scale)), max(1, int(height * scale))), Image.LANCZOS)

    def _score_ocr_text(self, text: str) -> float:
        cleaned = self._normalize_pdf_text(text)
        if not cleaned:
            return float("-inf")

        chinese = len(re.findall(r"[\u4e00-\u9fff]", cleaned))
        letters = len(re.findall(r"[A-Za-z]", cleaned))
        digits = len(re.findall(r"\d", cleaned))
        lines = len([line for line in cleaned.splitlines() if line.strip()])
        weird = cleaned.count("�")
        noise_blocks = len(re.findall(r"[_=~\-]{4,}", cleaned))
        return chinese * 2.0 + letters + digits * 0.8 + lines * 3.0 - weird * 15.0 - noise_blocks * 10.0

    def _get_ocr_backend(self, prefer_fast: bool = False):
        cache_key = "_ocr_backend_fast" if prefer_fast else "_ocr_backend"
        backend = getattr(self, cache_key, None)
        if backend is not None:
            return backend

        ocr_errors = []

        def load_tesseract():
            pytesseract = import_with_base_fallback('pytesseract')
            prepared = prepare_pytesseract()
            return {
                "name": "Tesseract",
                "kind": "pytesseract",
                "engine": pytesseract,
                "command": prepared["command"],
                "config": prepared.get("config", ""),
                "lang": prepared.get("lang", "chi_sim+eng"),
                "tessdata_dir": prepared.get("tessdata_dir", ""),
            }

        def load_easyocr():
            easyocr = import_with_base_fallback('easyocr')
            reader = easyocr.Reader(['ch_sim', 'en'], verbose=False)
            return {"name": "EasyOCR", "kind": "easyocr", "engine": reader}

        loaders = [
            ("pytesseract", load_tesseract),
            ("easyocr", load_easyocr),
        ] if prefer_fast else [
            ("easyocr", load_easyocr),
            ("pytesseract", load_tesseract),
        ]

        for loader_name, loader in loaders:
            try:
                backend = loader()
                setattr(self, cache_key, backend)
                self._active_ocr_engine = backend["name"]
                return backend
            except Exception as exc:
                ocr_errors.append(f"{loader_name}: {exc}")

        detail = "；".join(ocr_errors) if ocr_errors else "未检测到可用 OCR 引擎"
        raise RuntimeError(
            "PDF 解析失败：无法初始化 OCR 引擎。"
            "请安装 `easyocr`，或安装 `pytesseract` + Tesseract OCR 并确保可执行程序可用。"
            f"{detail}"
        )

    def _compute_stats(self, text, tables):
        return {
            "char_count": len(text),
            "paragraph_count": len([p for p in text.split('\n\n') if p.strip()]),
            "table_count": len(tables),
            "line_count": len(text.splitlines()),
        }


# ═══════════════════════════════════════════════════
#  Word 文档解析器
# ═══════════════════════════════════════════════════


class WordParser(BaseParser):
    """
    Word (.docx / .doc) 文档解析器
    - .docx: 直接用 python-docx 解析
    - .doc:  策略1: python-docx 直接打开（部分 .doc 实为 .docx 改名）
             策略2: 纯 Python OLE 解析（Word 97-2003 二进制格式，零外部依赖）
             策略3: soffice 转换（兜底，需安装 LibreOffice）
    依赖: pip install python-docx
    """
    supported_extensions = ('docx', 'doc')

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        start = time.time()
        result = self._make_result(file_path)
        ext = Path(file_path).suffix.lower()

        converted_path = None
        actual_path = file_path

        if ext == '.doc':
            # 策略1：python-docx 直接打开（.docx 改名为 .doc 的情况）
            try:
                from docx import Document as _T
                _T(file_path)
                actual_path = file_path
                logger.info(f".doc 实为 docx 格式，直接解析: {file_path}")
            except Exception:
                # 策略2：优先使用 soffice 转换（最可靠）
                converted_path = self._convert_doc_to_docx(file_path)
                if converted_path:
                    actual_path = converted_path
                    logger.info(f"soffice 转换成功: {file_path} -> {converted_path}")
                else:
                    # 策略3：soffice 不可用时，使用 OLE 解析兜底
                    try:
                        text = self._extract_doc_ole(file_path)
                        result.text_content = text
                        result.tables = []
                        result.metadata = {"原始格式": "DOC", "解析方式": "OLE内置解析器"}
                        result.processing_time_ms = (time.time() - start) * 1000
                        result.statistics = {
                            "char_count":      len(text),
                            "paragraph_count": len([p for p in text.split('\n') if p.strip()]),
                            "table_count":     0,
                        }
                        logger.info(f"Word(.doc)OLE解析完成: {file_path} — {len(text)} 字符")
                        return result
                    except Exception as ole_err:
                        result.error_msg = (
                            "无法处理 .doc 格式。请尝试：\n"
                            "① 用 Word/WPS 打开，另存为 .docx 后重新上传\n"
                            "② 安装 LibreOffice 后重启后端\n"
                            f"（OLE解析错误: {ole_err}）"
                        )
                        result.success = False
                        result.processing_time_ms = (time.time() - start) * 1000
                        return result
                    actual_path = converted_path

        try:
            from docx import Document
            doc = Document(actual_path)
            text_parts = []
            tables_data = []

            props = doc.core_properties
            result.metadata = {
                "title":    props.title or "",
                "author":   props.author or "",
                "created":  str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "",
                "revision": props.revision,
                "原始格式":  Path(file_path).suffix.upper().lstrip('.'),
            }

            # 收集表格段落 id，避免重复计入正文
            table_para_ids = set()
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            table_para_ids.add(id(para._element))

            for para in doc.paragraphs:
                if id(para._element) in table_para_ids:
                    continue
                if not para.text.strip():
                    continue
                style = para.style.name if para.style else ''
                if 'Heading' in style:
                    level = style.split()[-1] if style.split()[-1].isdigit() else '1'
                    text_parts.append(f"\n{'#' * int(level)} {para.text}")
                else:
                    text_parts.append(para.text)

            for table in doc.tables:
                rows = []
                seen_rows = set()
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    key = tuple(cells)
                    if key not in seen_rows:
                        seen_rows.add(key)
                        rows.append(cells)
                if rows:
                    tables_data.append(rows)

            result.text_content = "\n".join(text_parts)
            result.tables = tables_data

        except ImportError:
            result.error_msg = "缺少依赖: pip install python-docx"
            result.success = False
        except Exception as e:
            result.error_msg = str(e)
            result.success = False
        finally:
            if converted_path and os.path.exists(converted_path):
                try:
                    os.remove(converted_path)
                except Exception:
                    pass

        result.processing_time_ms = (time.time() - start) * 1000
        result.statistics = {
            "char_count":      len(result.text_content),
            "paragraph_count": len([p for p in result.text_content.split('\n') if p.strip()]),
            "table_count":     len(result.tables),
        }
        logger.info(f"Word解析完成: {file_path} — {len(result.tables)} 个表格，{len(result.text_content)} 字符")
        return result

    # ── 纯 Python OLE2 .doc 文本提取器 ──────────────────────────────
    @staticmethod
    def _extract_doc_ole(file_path: str) -> str:
        """
        无需任何第三方库，直接解析 Word 97-2003 OLE2 二进制格式，提取文本。
        优先按照 Word piece table 提取正文，避免直接解码整个 WordDocument 流导致乱码。
        """
        import struct, re

        with open(file_path, 'rb') as f:
            data = f.read()

        # 验证 OLE2 魔数
        if data[:8] != b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
            raise ValueError('非 OLE2 格式')

        sector_size      = 1 << struct.unpack_from('<H', data, 0x1e)[0]
        mini_sector_size = 1 << struct.unpack_from('<H', data, 0x20)[0]
        first_dir_sector = struct.unpack_from('<I', data, 0x30)[0]
        mini_stream_cutoff = struct.unpack_from('<I', data, 0x38)[0]
        first_minifat    = struct.unpack_from('<I', data, 0x3c)[0]

        FREESECT   = 0xFFFFFFFF
        ENDOFCHAIN = 0xFFFFFFFE

        def read_sector(sid):
            off = (sid + 1) * sector_size
            return data[off: off + sector_size]

        # 构建 FAT
        difat_sids = list(struct.unpack_from('<109I', data, 0x4c))
        fat_sids = [s for s in difat_sids if s not in (FREESECT, ENDOFCHAIN)]
        # 额外 DIFAT 扇区
        first_difat = struct.unpack_from('<I', data, 0x44)[0]
        if first_difat not in (FREESECT, ENDOFCHAIN):
            sid = first_difat
            while sid not in (FREESECT, ENDOFCHAIN):
                sec = read_sector(sid)
                entries = struct.unpack_from(f'<{sector_size // 4}I', sec)
                fat_sids.extend(s for s in entries[:-1] if s not in (FREESECT, ENDOFCHAIN))
                sid = entries[-1]

        fat = []
        for sid in fat_sids:
            sec = read_sector(sid)
            fat.extend(struct.unpack_from(f'<{sector_size // 4}I', sec))

        def follow_chain(start):
            chain, sid, seen = [], start, set()
            while sid not in (FREESECT, ENDOFCHAIN) and sid not in seen and sid < len(fat):
                seen.add(sid); chain.append(sid); sid = fat[sid]
            return chain

        def read_stream(start, size=None):
            raw = b''.join(read_sector(s) for s in follow_chain(start))
            return raw[:size] if size is not None else raw

        # 目录
        dir_data = read_stream(first_dir_sector)
        # 根目录 → 获取 mini stream
        root = dir_data[:128]
        root_start = struct.unpack_from('<I', root, 0x74)[0]
        mini_stream = read_stream(root_start) if root_start not in (FREESECT, ENDOFCHAIN) else b''

        # Mini FAT
        minifat = []
        if first_minifat not in (FREESECT, ENDOFCHAIN):
            mf_data = read_stream(first_minifat)
            minifat = list(struct.unpack_from(f'<{len(mf_data) // 4}I', mf_data))

        def follow_mini_chain(start):
            chain, sid, seen = [], start, set()
            while sid not in (FREESECT, ENDOFCHAIN) and sid not in seen and sid < len(minifat):
                seen.add(sid); chain.append(sid); sid = minifat[sid]
            return chain

        def read_mini(start, size):
            raw = b''.join(mini_stream[s * mini_sector_size:(s + 1) * mini_sector_size]
                           for s in follow_mini_chain(start))
            return raw[:size]

        # 解析目录条目，找 WordDocument 流
        entries = {}
        for i in range(len(dir_data) // 128):
            entry = dir_data[i * 128:(i + 1) * 128]
            name_len = struct.unpack_from('<H', entry, 0x40)[0]
            if name_len < 2: continue
            name = entry[:name_len - 2].decode('utf-16-le', errors='ignore')
            start_sid = struct.unpack_from('<I', entry, 0x74)[0]
            size      = struct.unpack_from('<I', entry, 0x78)[0]
            entries[name] = (start_sid, size)

        if 'WordDocument' not in entries:
            raise ValueError('WordDocument 流不存在')

        def read_entry_stream(name: str) -> bytes:
            if name not in entries:
                raise ValueError(f'{name} 流不存在')
            start_sid, size = entries[name]
            if size < mini_stream_cutoff and minifat:
                return read_mini(start_sid, size)
            return read_stream(start_sid, size)

        def decode_single_byte_piece(raw_bytes: bytes, language_id: int) -> str:
            codec = {
                0x0804: 'gbk',       # 简体中文（中国）
                0x1004: 'gbk',       # 简体中文（新加坡）
                0x0404: 'big5',      # 繁体中文（台湾）
                0x0C04: 'big5',      # 繁体中文（香港）
                0x0411: 'shift_jis', # 日文
                0x0412: 'cp949',     # 韩文
            }.get(language_id, 'cp1252')

            for enc in (codec, 'cp1252', 'latin-1'):
                try:
                    return raw_bytes.decode(enc, errors='ignore')
                except LookupError:
                    continue
            return raw_bytes.decode('latin-1', errors='ignore')

        word_data = read_entry_stream('WordDocument')
        fib_flags = struct.unpack_from('<H', word_data, 0x0A)[0]
        language_id = struct.unpack_from('<H', word_data, 0x06)[0]
        table_stream_name = '1Table' if (fib_flags & 0x0200) else '0Table'
        table_data = read_entry_stream(table_stream_name)

        fib_off = 0x20
        csw = struct.unpack_from('<H', word_data, fib_off)[0]
        fib_off += 2
        fib_off += csw * 2

        cslw = struct.unpack_from('<H', word_data, fib_off)[0]
        fib_off += 2
        fib_lw = word_data[fib_off:fib_off + cslw * 4]
        if len(fib_lw) < 16:
            raise ValueError('DOC FIB 结构损坏，无法读取正文长度')
        ccp_text = struct.unpack_from('<I', fib_lw, 12)[0]
        fib_off += cslw * 4

        cb_rg_fc_lcb = struct.unpack_from('<H', word_data, fib_off)[0]
        fib_off += 2
        if cb_rg_fc_lcb < 34:
            raise ValueError('DOC FIB 结构不完整，无法定位 CLX')

        # FibRgFcLcb97 中第 34 对字段是 fcClx/lcbClx
        fc_clx, lcb_clx = struct.unpack_from('<II', word_data, fib_off + 33 * 8)
        clx = table_data[fc_clx:fc_clx + lcb_clx]
        if not clx:
            raise ValueError('CLX 为空，无法解析 DOC 文本')

        pos = 0
        while pos < len(clx) and clx[pos] == 0x01:
            if pos + 3 > len(clx):
                raise ValueError('CLX PRC 结构损坏')
            cb_grpprl = struct.unpack_from('<H', clx, pos + 1)[0]
            pos += 3 + cb_grpprl

        if pos >= len(clx) or clx[pos] != 0x02:
            raise ValueError('未找到 DOC piece table')

        piece_table_size = struct.unpack_from('<I', clx, pos + 1)[0]
        piece_table = clx[pos + 5:pos + 5 + piece_table_size]
        if piece_table_size < 4 or len(piece_table) < piece_table_size:
            raise ValueError('DOC piece table 损坏')

        piece_count = (piece_table_size - 4) // 12
        if piece_count <= 0:
            raise ValueError('DOC piece table 为空')

        cp_list = list(struct.unpack_from(f'<{piece_count + 1}I', piece_table, 0))
        text_parts = []

        for i in range(piece_count):
            cp_start = cp_list[i]
            cp_end = min(cp_list[i + 1], ccp_text)
            if cp_start >= ccp_text or cp_end <= cp_start:
                continue

            pcd_off = 4 * (piece_count + 1) + i * 8
            fc_raw = struct.unpack_from('<I', piece_table, pcd_off + 2)[0]
            is_single_byte = bool(fc_raw & 0x40000000)
            fc_value = fc_raw & 0x3FFFFFFF
            char_count = cp_end - cp_start

            if is_single_byte:
                byte_offset = fc_value // 2
                chunk = word_data[byte_offset:byte_offset + char_count]
                piece_text = decode_single_byte_piece(chunk, language_id)
            else:
                byte_offset = fc_value
                chunk = word_data[byte_offset:byte_offset + char_count * 2]
                piece_text = chunk.decode('utf-16-le', errors='ignore')

            text_parts.append(piece_text)

        raw = ''.join(text_parts)
        raw = (
            raw.replace('\r', '\n')
               .replace('\x0b', '\n')
               .replace('\x0c', '\n')
               .replace('\x07', '\t')
               .replace('\xa0', ' ')
               .replace('\x13', '')
               .replace('\x14', '')
               .replace('\x15', '')
        )

        # 清理常见控制词和不可见字符
        clean = re.sub(r'[\x00-\x08\x0e-\x1f\x7f-\x9f]+', '', raw)
        clean = re.sub(r'HYPERLINK|PAGEREF|MERGEFORMAT|TOC|EMBED|Visio\.Drawing\.\d+|Toc\d+|FORMTEXT|MACROBUTTON', ' ', clean)
        clean = re.sub(r'[ \t]{3,}', '  ', clean)
        clean = re.sub(r'\n[ \t]+', '\n', clean)
        clean = re.sub(r'\n{3,}', '\n\n', clean)

        # 按行做一次轻量过滤，剔除残留的异常碎片
        lines = clean.split('\n')
        valid_lines = []
        for line in lines:
            line = line.strip()
            if not line:
                continue

            cn_chars = len(re.findall(r'[\u4e00-\u9fff]', line))
            en_chars = len(re.findall(r'[a-zA-Z]', line))
            digit_chars = len(re.findall(r'\d', line))
            text_chars = cn_chars + en_chars + digit_chars
            total_chars = len(line)

            if len(set(line)) < len(line) * 0.25 and len(line) > 10:
                continue

            ratio = text_chars / total_chars if total_chars > 0 else 0
            if ratio > 0.5 or (total_chars > 10 and ratio > 0.35):
                valid_lines.append(line)

        text = '\n'.join(valid_lines).strip()
        if not text:
            raise ValueError('DOC 文本为空或无法识别')
        return text

    # ── soffice 兜底（仅在 OLE 解析失败时使用）─────────────────────
    def _convert_doc_to_docx(self, file_path: str) -> Optional[str]:
        import subprocess, tempfile, shutil, platform
        system = platform.system()
        candidates = {
            'Windows': [
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            ],
            'Darwin': [
                '/Applications/LibreOffice.app/Contents/MacOS/soffice',
                '/opt/homebrew/bin/soffice',
            ],
        }.get(system, [
            '/usr/bin/soffice', '/usr/bin/libreoffice',
            '/usr/lib/libreoffice/program/soffice',
        ])

        soffice = next((c for c in candidates if os.path.isfile(c) and os.access(c, os.X_OK)), None)
        if not soffice:
            env_path = '/usr/bin:/usr/local/bin:' + os.environ.get('PATH', '')
            soffice = shutil.which('soffice', path=env_path) or shutil.which('libreoffice', path=env_path)
        if not soffice:
            return None

        try:
            abs_file     = os.path.abspath(file_path)
            tmp_dir      = tempfile.mkdtemp()
            user_profile = tempfile.mkdtemp()
            env = os.environ.copy()
            env['PATH'] = '/usr/bin:/usr/local/bin:' + env.get('PATH', '')
            cmd = [soffice,
                   f'-env:UserInstallation=file://{user_profile}',
                   '--headless', '--norestore', '--nofirststartwizard',
                   '--convert-to', 'docx', '--outdir', tmp_dir, abs_file]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=90, env=env)
            if r.returncode != 0:
                logger.error(f"soffice 失败: {r.stderr[:200]}")
                return None
            stem = Path(abs_file).stem
            out = Path(tmp_dir) / f"{stem}.docx"
            if out.exists():
                return str(out)
            for f in Path(tmp_dir).glob("*.docx"):
                return str(f)
            return None
        except Exception as e:
            logger.error(f"soffice 异常: {e}")
            return None


class ExcelParser(BaseParser):
    """
    Excel (.xlsx/.xls) 解析器
    依赖: pip install openpyxl
    """
    supported_extensions = ('xlsx', 'xls')

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        start = time.time()
        result = self._make_result(file_path)
        ext = Path(file_path).suffix.lower()

        try:
            import openpyxl
            try:
                wb = openpyxl.load_workbook(file_path, data_only=True)
            except Exception:
                if ext == '.xls':
                    with open(file_path, 'rb') as stream:
                        wb = openpyxl.load_workbook(stream, data_only=True)
                else:
                    raise

            result.metadata = {
                "sheet_names": wb.sheetnames,
                "sheet_count": len(wb.sheetnames),
            }

            text_parts = []
            all_tables = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n[工作表: {sheet_name}]")

                rows = []
                for row in ws.iter_rows(values_only=True):
                    clean_row = [str(c) if c is not None else "" for c in row]
                    # 跳过全空行
                    if any(c for c in clean_row):
                        rows.append(clean_row)
                        text_parts.append("\t".join(clean_row))

                if rows:
                    all_tables.append(rows)

            result.text_content = "\n".join(text_parts)
            result.tables = all_tables

            # 数值统计
            nums = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if isinstance(cell, (int, float)):
                            nums.append(cell)

            if nums:
                result.statistics["numeric_count"] = len(nums)
                result.statistics["sum"] = sum(nums)
                result.statistics["avg"] = sum(nums) / len(nums)
                result.statistics["max"] = max(nums)
                result.statistics["min"] = min(nums)

        except ImportError:
            result.error_msg = "缺少依赖: pip install openpyxl"
            result.success = False
        except Exception as e:
            result.error_msg = str(e)
            result.success = False

        result.processing_time_ms = (time.time() - start) * 1000
        logger.info(f"Excel解析完成: {file_path} — {len(result.tables)} 个Sheet")
        return result


# ═══════════════════════════════════════════════════
#  PPT 解析器
# ═══════════════════════════════════════════════════

class PPTXParser(BaseParser):
    """
    PowerPoint (.pptx) 解析器
    依赖: pip install python-pptx
    """
    supported_extensions = ('pptx',)

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        start = time.time()
        result = self._make_result(file_path)

        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            result.metadata = {
                "slide_count": len(prs.slides),
                "slide_width":  prs.slide_width.cm if prs.slide_width else None,
                "slide_height": prs.slide_height.cm if prs.slide_height else None,
            }

            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))

                # 提取备注
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_parts.append(f"  [备注] {notes}")

            result.text_content = "\n\n".join(text_parts)

        except ImportError:
            result.error_msg = "缺少依赖: pip install python-pptx"
            result.success = False
        except Exception as e:
            result.error_msg = str(e)
            result.success = False

        result.processing_time_ms = (time.time() - start) * 1000
        result.statistics = {
            "char_count": len(result.text_content),
            "slide_count": result.metadata.get("slide_count", 0),
        }
        logger.info(f"PPTX解析完成: {file_path} — {result.metadata.get('slide_count',0)} 张幻灯片")
        return result


# ═══════════════════════════════════════════════════
#  TXT / Markdown / JSON 解析器
# ═══════════════════════════════════════════════════

class TextParser(BaseParser):
    """纯文本解析器，支持 txt / md / json / csv"""
    supported_extensions = ('txt', 'md', 'json', 'csv')

    def parse(self, file_path: str, **kwargs) -> ExtractionResult:
        start = time.time()
        result = self._make_result(file_path)
        ext = Path(file_path).suffix.lower()

        try:
            # 自动检测编码
            content = self._read_with_encoding(file_path)

            if ext == '.json':
                obj = json.loads(content)
                result.text_content = json.dumps(obj, ensure_ascii=False, indent=2)
                result.metadata = {"keys": list(obj.keys()) if isinstance(obj, dict) else [], "type": type(obj).__name__}

            elif ext == '.csv':
                import csv, io
                reader = csv.reader(io.StringIO(content))
                rows = list(reader)
                result.tables = [rows]
                result.text_content = content
                result.metadata = {"rows": len(rows), "cols": len(rows[0]) if rows else 0}

            else:
                result.text_content = content

            result.metadata.update({
                "encoding": "UTF-8",
                "file_size_bytes": os.path.getsize(file_path),
                "line_count": len(content.splitlines()),
            })

        except Exception as e:
            result.error_msg = str(e)
            result.success = False

        result.processing_time_ms = (time.time() - start) * 1000
        result.statistics = {
            "char_count": len(result.text_content),
            "word_count": len(result.text_content.split()),
            "line_count": len(result.text_content.splitlines()),
        }
        return result

    def _read_with_encoding(self, path: str) -> str:
        for enc in ('utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1'):
            try:
                with open(path, 'r', encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, LookupError):
                continue
        raise ValueError("无法识别文件编码")


# ═══════════════════════════════════════════════════
#  关键词提取 & 摘要工具
# ═══════════════════════════════════════════════════

class TextAnalyzer:
    """文本分析工具：关键词提取、词频统计、自动摘要"""

    # 中文停用词（精简版）
    STOP_WORDS = set(
        "的了是在和有与一这个为中对上以到其都也被从及等将但是不已"
        "就很再次并该这些那些之所以因为由于如果只有还有另外除了"
    )

    def extract_keywords(self, text: str, top_n: int = 15) -> list[tuple[str, int]]:
        """基于词频的关键词提取（简版，无需分词库）"""
        # 提取2~6字的中文词组（简单 n-gram）
        chinese_chars = re.findall(r'[\u4e00-\u9fa5]+', text)
        candidates = {}
        for word in chinese_chars:
            for length in range(2, 5):
                for start in range(len(word) - length + 1):
                    term = word[start:start+length]
                    if not any(c in self.STOP_WORDS for c in term):
                        candidates[term] = candidates.get(term, 0) + 1
        # 过滤低频
        filtered = {k: v for k, v in candidates.items() if v >= 2}
        return sorted(filtered.items(), key=lambda x: -x[1])[:top_n]

    def word_frequency(self, text: str) -> dict:
        """词频统计"""
        words = re.findall(r'[\u4e00-\u9fa5a-zA-Z]{2,}', text)
        freq = {}
        for w in words:
            if w not in self.STOP_WORDS:
                freq[w] = freq.get(w, 0) + 1
        return dict(sorted(freq.items(), key=lambda x: -x[1])[:50])

    def extract_summary(self, text: str, max_sentences: int = 5) -> str:
        """基于TF权重的抽取式摘要"""
        sentences = re.split(r'[。！？\n]', text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        if not sentences:
            return text[:200]

        keywords = dict(self.extract_keywords(text, top_n=20))

        def score(s):
            return sum(keywords.get(s[i:i+2], 0) for i in range(len(s)-1))

        scored = sorted(enumerate(sentences), key=lambda x: score(x[1]), reverse=True)
        top_idx = sorted([i for i, _ in scored[:max_sentences]])
        return '。'.join(sentences[i] for i in top_idx) + '。'


# ═══════════════════════════════════════════════════
#  输出格式化
# ═══════════════════════════════════════════════════

class OutputFormatter:
    """将提取结果转换为各种输出格式"""

    @staticmethod
    def to_txt(result: ExtractionResult) -> str:
        lines = [
            f"╔{'═'*60}",
            f"║ 文件: {Path(result.file_path).name}",
            f"║ 格式: {result.file_type.upper()}  |  耗时: {result.processing_time_ms:.1f}ms",
            f"╚{'═'*60}",
            "",
            "【文本内容】",
            result.text_content,
        ]
        if result.tables:
            lines.append(f"\n【表格数据】共 {len(result.tables)} 个表格")
            for i, table in enumerate(result.tables):
                lines.append(f"\n表格 {i+1}:")
                for row in table:
                    lines.append("  " + " | ".join(str(c) for c in row))
        if result.metadata:
            lines.append("\n【元数据】")
            for k, v in result.metadata.items():
                lines.append(f"  {k}: {v}")
        return "\n".join(lines)

    @staticmethod
    def to_json(result: ExtractionResult) -> str:
        return json.dumps(result.to_dict(), ensure_ascii=False, indent=2)

    @staticmethod
    def to_markdown(result: ExtractionResult) -> str:
        name = Path(result.file_path).name
        lines = [
            f"# 文档提取报告: {name}",
            "",
            f"> 格式: `{result.file_type.upper()}` | 处理耗时: `{result.processing_time_ms:.1f}ms`",
            "",
            "## 文本内容",
            "",
            result.text_content,
        ]
        if result.tables:
            lines.append("\n## 表格数据")
            for i, table in enumerate(result.tables):
                lines.append(f"\n### 表格 {i+1}")
                if table:
                    # Markdown table
                    header = table[0]
                    lines.append("| " + " | ".join(str(h) for h in header) + " |")
                    lines.append("|" + "|".join(["---"]*len(header)) + "|")
                    for row in table[1:]:
                        lines.append("| " + " | ".join(str(c) for c in row) + " |")
        return "\n".join(lines)

    @staticmethod
    def to_csv(result: ExtractionResult) -> str:
        import csv, io
        output = io.StringIO()
        writer = csv.writer(output)
        for table in result.tables:
            for row in table:
                writer.writerow(row)
            writer.writerow([])
        return output.getvalue()


# ═══════════════════════════════════════════════════
#  主处理器（门面模式）
# ═══════════════════════════════════════════════════

class DocFlowProcessor:
    """
    文档处理器主类 —— 门面模式
    自动选择合适的解析器处理文件
    """

    def __init__(self):
        self.parsers: list[BaseParser] = [
            PDFParser(),
            WordParser(),
            ExcelParser(),
            PPTXParser(),
            TextParser(),
        ]
        self.analyzer  = TextAnalyzer()
        self.formatter = OutputFormatter()
        logger.info("DocFlow 初始化完成，已注册 %d 个解析器", len(self.parsers))

    def process(
        self,
        file_path: str,
        extract_keywords: bool = True,
        output_format: str = "txt",   # "txt" | "json" | "markdown" | "csv"
        pdf_mode: str = "balanced",
        progress_callback=None,
        cancel_callback=None,
    ) -> dict:
        """处理单个文件，返回结构化结果字典"""
        def emit(progress_pct: float, stage: str, message: str = "", **extra) -> None:
            if callable(progress_callback):
                progress_callback(
                    progress_pct=max(0.0, min(float(progress_pct), 100.0)),
                    stage=stage,
                    message=message,
                    **extra,
                )

        path = Path(file_path)
        if not path.exists():
            return {"success": False, "error": f"文件不存在: {file_path}"}

        _ensure_not_cancelled(cancel_callback)
        emit(2, "prepare", f"正在检查文件：{path.name}")
        # 找到合适的解析器
        parser = self._find_parser(file_path)
        if not parser:
            return {"success": False, "error": f"不支持的文件格式: {path.suffix}"}

        emit(8, "parser_ready", f"已识别文件类型：{path.suffix.lower() or parser.__class__.__name__}")
        logger.info("开始处理: %s", path.name)
        if isinstance(parser, PDFParser):
            def parser_progress(progress_pct: float, stage: str = "", message: str = "", **extra) -> None:
                mapped = 10 + max(0.0, min(float(progress_pct), 100.0)) * 0.74
                emit(mapped, stage or "parsing", message or "正在解析文档", **extra)

            result = parser.parse(file_path, pdf_mode=pdf_mode, progress_callback=parser_progress, cancel_callback=cancel_callback)
        else:
            emit(18, "parsing", "正在解析文档内容")
            result = parser.parse(file_path, pdf_mode=pdf_mode, cancel_callback=cancel_callback)
            emit(78, "parsed", "文档解析完成，正在整理内容")

        # 关键词 & 摘要
        _ensure_not_cancelled(cancel_callback)
        if extract_keywords and result.text_content:
            emit(88, "analyzing", "正在提取关键词和自动摘要")
            result.statistics["keywords"] = self.analyzer.extract_keywords(result.text_content)
            result.statistics["summary"]  = self.analyzer.extract_summary(result.text_content)
        else:
            emit(88, "analyzing", "正在整理结构化结果")

        # 格式化输出
        fmt_map = {
            "txt":      self.formatter.to_txt,
            "json":     self.formatter.to_json,
            "markdown": self.formatter.to_markdown,
            "csv":      self.formatter.to_csv,
        }
        fmt_fn  = fmt_map.get(output_format, self.formatter.to_txt)
        _ensure_not_cancelled(cancel_callback)
        emit(96, "formatting", f"正在生成 {output_format.upper()} 输出")
        formatted = fmt_fn(result)
        emit(100, "done", "处理完成")

        return {
            "success":        result.success,
            "file":           path.name,
            "format":         result.file_type,
            "text":           result.text_content,
            "tables":         result.tables,
            "metadata":       result.metadata,
            "statistics":     result.statistics,
            "processing_ms":  result.processing_time_ms,
            "formatted_output": formatted,
            "error":          result.error_msg,
        }

    def process_batch(self, file_paths: list[str], **kwargs) -> list[dict]:
        """批量处理多个文件"""
        results = []
        for i, fp in enumerate(file_paths):
            logger.info("[%d/%d] 处理: %s", i+1, len(file_paths), Path(fp).name)
            results.append(self.process(fp, **kwargs))
        return results

    def _find_parser(self, file_path: str) -> Optional[BaseParser]:
        for parser in self.parsers:
            if parser.can_parse(file_path):
                return parser
        return None


# ═══════════════════════════════════════════════════
#  命令行入口
# ═══════════════════════════════════════════════════

def main():
    import argparse

    parser = argparse.ArgumentParser(
        prog="docflow",
        description="DocFlow — 多格式文档自动化处理工具"
    )
    parser.add_argument("files", nargs="+", help="待处理的文件路径")
    parser.add_argument(
        "-f", "--format",
        choices=["txt", "json", "markdown", "csv"],
        default="txt",
        help="输出格式 (默认: txt)"
    )
    parser.add_argument(
        "-o", "--output",
        help="结果输出目录（默认打印到终端）"
    )
    parser.add_argument(
        "--no-keywords",
        action="store_true",
        help="跳过关键词提取"
    )
    parser.add_argument(
        "--pdf-mode",
        choices=["accurate", "balanced", "fast"],
        default="balanced",
        help="PDF 解析模式（默认: balanced）",
    )

    args = parser.parse_args()

    processor = DocFlowProcessor()

    for file_path in args.files:
        result = processor.process(
            file_path,
            extract_keywords=not args.no_keywords,
            output_format=args.format,
            pdf_mode=args.pdf_mode,
        )

        if not result["success"]:
            print(f"[错误] {result['file']}: {result['error']}")
            continue

        output_text = result["formatted_output"]

        if args.output:
            out_dir = Path(args.output)
            out_dir.mkdir(parents=True, exist_ok=True)
            stem = Path(file_path).stem
            ext_map = {"txt":"txt","json":"json","markdown":"md","csv":"csv"}
            out_file = out_dir / f"{stem}_extracted.{ext_map[args.format]}"
            out_file.write_text(output_text, encoding="utf-8")
            print(f"✓ {result['file']} → {out_file}")
        else:
            print(output_text)
            print(f"\n处理耗时: {result['processing_ms']:.1f}ms")

        # 打印统计
        stats = result.get("statistics", {})
        if stats.get("keywords"):
            kws = ", ".join(f"{w}({c})" for w, c in stats["keywords"][:5])
            print(f"关键词: {kws}")
        if stats.get("summary"):
            print(f"摘要: {stats['summary'][:100]}...")


if __name__ == "__main__":
    main()
